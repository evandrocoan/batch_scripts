import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pathlib import Path
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor

# Default configuration
DEFAULT_GLOW_COLOR = "#FFFFF0"
DEFAULT_GLOW_SIZE = 18
DEFAULT_TEXT_COLOR = "#010101"

class SlideFixerGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Slide Fixer - OBS Chroma Key")
        self.root.geometry("600x500")
        self.root.resizable(False, False)
        
        self.selected_file = None
        
        # Create UI
        self.create_widgets()
    
    def create_widgets(self):
        # Title
        title_label = tk.Label(
            self.root, 
            text="PowerPoint Slide Fixer for OBS",
            font=("Arial", 16, "bold")
        )
        title_label.pack(pady=20)
        
        # File selection frame
        file_frame = tk.Frame(self.root)
        file_frame.pack(pady=10, padx=20, fill="x")
        
        tk.Label(file_frame, text="Select PowerPoint File:").pack(anchor="w")
        
        self.file_entry = tk.Entry(file_frame, width=60, state="readonly")
        self.file_entry.pack(side="left", padx=(0, 10), fill="x", expand=True)
        
        browse_btn = tk.Button(file_frame, text="Browse...", command=self.browse_file, width=10)
        browse_btn.pack(side="left")
        
        # Configuration frame
        config_frame = tk.LabelFrame(self.root, text="Configuration", padx=20, pady=15)
        config_frame.pack(pady=20, padx=20, fill="both")
        
        # Glow color
        color_frame = tk.Frame(config_frame)
        color_frame.pack(fill="x", pady=5)
        tk.Label(color_frame, text="Glow Color (hex):", width=20, anchor="w").pack(side="left")
        self.color_entry = tk.Entry(color_frame, width=15)
        self.color_entry.insert(0, DEFAULT_GLOW_COLOR)
        self.color_entry.pack(side="left", padx=5)
        
        # Glow size
        size_frame = tk.Frame(config_frame)
        size_frame.pack(fill="x", pady=5)
        tk.Label(size_frame, text="Glow Size (points):", width=20, anchor="w").pack(side="left")
        self.size_spinbox = tk.Spinbox(size_frame, from_=10, to=50, width=13)
        self.size_spinbox.delete(0, "end")
        self.size_spinbox.insert(0, DEFAULT_GLOW_SIZE)
        self.size_spinbox.pack(side="left", padx=5)
        
        # Text color
        text_color_frame = tk.Frame(config_frame)
        text_color_frame.pack(fill="x", pady=5)
        tk.Label(text_color_frame, text="Text Color (hex):", width=20, anchor="w").pack(side="left")
        self.text_color_entry = tk.Entry(text_color_frame, width=15)
        self.text_color_entry.insert(0, DEFAULT_TEXT_COLOR)
        self.text_color_entry.pack(side="left", padx=5)
        
        # Progress bar
        self.progress_frame = tk.Frame(self.root)
        self.progress_frame.pack(pady=10, padx=20, fill="x")
        
        self.progress_label = tk.Label(self.progress_frame, text="")
        self.progress_label.pack()
        
        self.progress_bar = ttk.Progressbar(
            self.progress_frame, 
            mode='indeterminate', 
            length=500
        )
        
        # Process button
        button_frame = tk.Frame(self.root)
        button_frame.pack(pady=20, fill="x")
        
        self.process_btn = tk.Button(
            button_frame,
            text="✓ Process Presentation",
            command=self.process_file,
            width=25,
            height=2,
            bg="#4CAF50",
            fg="white",
            font=("Arial", 14, "bold"),
            state="disabled",
            cursor="hand2"
        )
        self.process_btn.pack()
    
    def browse_file(self):
        filename = filedialog.askopenfilename(
            title="Select PowerPoint Presentation",
            filetypes=[
                ("PowerPoint files", "*.pptx"),
                ("All files", "*.*")
            ]
        )
        
        if filename:
            self.selected_file = filename
            self.file_entry.config(state="normal")
            self.file_entry.delete(0, tk.END)
            self.file_entry.insert(0, filename)
            self.file_entry.config(state="readonly")
            self.process_btn.config(state="normal")
    
    def apply_solid_glow_to_run(self, run, color_hex, size_pt):
        """Applies a SOLID glow effect to create a highlighter background effect."""
        # Strip '#' from color if present
        color_hex = color_hex.lstrip('#')
        
        # Get the Run Properties (rPr) element or create it
        rPr = run._r.get_or_add_rPr()
        
        # Calculate radius in EMUs (1 point = 12700 EMUs)
        radius_emu = int(size_pt * 12700)
        
        # Create the XML for the effect list with glow
        effectlst_xml = f'''<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:glow rad="{radius_emu}">
                <a:srgbClr val="{color_hex}">
                    <a:alpha val="100000"/>
                </a:srgbClr>
            </a:glow>
        </a:effectLst>'''
        
        # Remove any existing effectLst or standalone glow
        existing_effectlst = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
        if existing_effectlst is not None:
            rPr.remove(existing_effectlst)
        
        existing_glow = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}glow")
        if existing_glow is not None:
            rPr.remove(existing_glow)
        
        # Parse the new XML and inject it into the run properties
        effectlst_element = parse_xml(effectlst_xml)
        
        # Insert after solidFill if it exists, otherwise at the beginning
        solid_fill = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
        if solid_fill is not None:
            insert_index = list(rPr).index(solid_fill) + 1
            rPr.insert(insert_index, effectlst_element)
        else:
            rPr.insert(0, effectlst_element)
    
    def process_file(self):
        if not self.selected_file:
            messagebox.showerror("Error", "Please select a file first!")
            return
        
        try:
            # Get configuration values
            glow_color = self.color_entry.get()
            glow_size = int(self.size_spinbox.get())
            text_color = self.text_color_entry.get().lstrip('#')
            
            # Validate colors
            if not glow_color or len(glow_color.lstrip('#')) != 6:
                messagebox.showerror("Error", "Invalid glow color! Use format: #RRGGBB")
                return
            
            if not text_color or len(text_color) != 6:
                messagebox.showerror("Error", "Invalid text color! Use format: #RRGGBB")
                return
            
            # Disable button and show progress
            self.process_btn.config(state="disabled")
            self.progress_label.config(text="Processing presentation...")
            self.progress_bar.pack(pady=5)
            self.progress_bar.start(10)
            self.root.update()
            
            # Generate output filename
            input_path = Path(self.selected_file)
            output_path = input_path.parent / f"{input_path.stem}_obs_fixed{input_path.suffix}"
            
            # Open presentation
            prs = Presentation(self.selected_file)
            count = 0
            
            # Process each slide
            for slide in prs.slides:
                # Check if slide has any text
                has_text = False
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        has_text = True
                        break
                
                # Set slide background: white if has text, black if no text
                background = slide.background
                fill = background.fill
                fill.solid()
                if has_text:
                    fill.fore_color.rgb = RGBColor(255, 255, 255)
                else:
                    fill.fore_color.rgb = RGBColor(0, 0, 0)
                
                # Process all text shapes
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    text_processed = False
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                # Apply the SOLID glow
                                self.apply_solid_glow_to_run(run, glow_color, glow_size)
                                
                                # Set text color
                                run.font.color.rgb = RGBColor(
                                    int(text_color[0:2], 16),
                                    int(text_color[2:4], 16),
                                    int(text_color[4:6], 16)
                                )
                                text_processed = True
                    
                    if text_processed:
                        count += 1
            
            # Save the presentation
            prs.save(str(output_path))
            
            # Stop progress bar
            self.progress_bar.stop()
            self.progress_bar.pack_forget()
            self.progress_label.config(text="")
            self.process_btn.config(state="normal")
            
            # Show success message
            messagebox.showinfo(
                "Success",
                f"Processed {count} text shapes!\n\nSaved to:\n{output_path.name}"
            )
            
        except Exception as e:
            self.progress_bar.stop()
            self.progress_bar.pack_forget()
            self.progress_label.config(text="")
            self.process_btn.config(state="normal")
            messagebox.showerror("Error", f"An error occurred:\n{str(e)}")

def main():
    root = tk.Tk()
    app = SlideFixerGUI(root)
    root.mainloop()

if __name__ == "__main__":
    main()
