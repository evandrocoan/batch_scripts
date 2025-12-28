from pptx import Presentation
from pptx.util import Pt
import sys

# Make a copy and process
import shutil
shutil.copy('Apresentação1Original.pptx', 'test_output.pptx')

prs = Presentation('test_output.pptx')

print('=== PROCESSING SLIDE BY SLIDE ===', flush=True)

from pptx.enum.text import PP_ALIGN

slide_width = prs.slide_width
slide_height = prs.slide_height

margin_percent = 0.05
spacing_pt = 10

margin_x = int(slide_width * margin_percent)
margin_y = int(slide_height * margin_percent)
spacing_emu = int(spacing_pt * 12700)

available_width = slide_width - (2 * margin_x)
available_height = slide_height - (2 * margin_y)
available_width_pt = available_width / 12700
available_height_pt = available_height / 12700

print(f'Slide size: {slide_width/914400:.1f}" x {slide_height/914400:.1f}"', flush=True)
print(f'Available: {available_width_pt:.0f}pt x {available_height_pt:.0f}pt', flush=True)

for slide_num, slide in enumerate(prs.slides, 1):
    print(f'\nProcessing slide {slide_num}...', flush=True)
    
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()
            lines = text.count('\n') + 1
            char_count = len(text)
            weight = max(1, lines + (char_count / 50))
            
            original_sizes = {}
            min_font_in_shape = float('inf')
            max_font_in_shape = 0
            font_name = 'Arial'
            
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(paragraph.runs):
                    if run.text.strip():
                        if run.font.size:
                            font_size = run.font.size.pt
                        else:
                            font_size = 12
                        
                        original_sizes[(para_idx, run_idx)] = font_size
                        min_font_in_shape = min(min_font_in_shape, font_size)
                        max_font_in_shape = max(max_font_in_shape, font_size)
                        
                        if run.font.name:
                            font_name = run.font.name
            
            if min_font_in_shape == float('inf'):
                min_font_in_shape = 12
            if max_font_in_shape == 0:
                max_font_in_shape = 12
            
            print(f'  Shape "{shape.name}": fonts {min_font_in_shape}-{max_font_in_shape}pt, {len(original_sizes)} runs', flush=True)
            
            text_shapes.append({
                'shape': shape, 
                'weight': weight, 
                'text': text,
                'original_sizes': original_sizes,
                'min_font': min_font_in_shape,
                'max_font': max_font_in_shape,
                'font_name': font_name
            })
    
    if not text_shapes:
        print('  (no text shapes)', flush=True)
        continue
    
    print(f'  Found {len(text_shapes)} text shapes', flush=True)
    
    # Just process first 5 slides for testing
    if slide_num > 5:
        print('  (skipping rest for testing)', flush=True)
        continue
    
    num_shapes = len(text_shapes)
    total_weight = sum(item['weight'] for item in text_shapes)
    total_spacing = spacing_emu * (num_shapes - 1) if num_shapes > 1 else 0
    height_for_boxes = available_height - total_spacing
    
    current_y = margin_y
    shape_layout = []
    
    for item in text_shapes:
        height_ratio = item['weight'] / total_weight
        height_per_box = int(height_for_boxes * height_ratio)
        min_height = int(available_height * 0.1)
        height_per_box = max(height_per_box, min_height)
        
        shape_layout.append({
            'item': item,
            'height': height_per_box,
            'height_pt': height_per_box / 12700,
            'y': current_y
        })
        current_y += height_per_box + spacing_emu
    
    # Calculate scale factors
    import fix_slides_for_obs_processor as processor
    
    for layout in shape_layout:
        item = layout['item']
        height_pt = layout['height_pt']
        original_sizes = item['original_sizes']
        max_font = item['max_font']
        font_name = item['font_name']
        text_content = item['text']
        
        print(f'    Calculating scale for "{item["shape"].name}" (max_font={max_font}, height={height_pt:.0f}pt)...', flush=True)
        
        if not original_sizes:
            layout['scale_factor'] = 1.0
            continue
        
        low_scale = 0.1
        high_scale = 20.0
        best_scale = 1.0
        margin_pt = 10
        
        iterations = 0
        for _ in range(20):
            iterations += 1
            mid_scale = (low_scale + high_scale) / 2
            max_scaled_font = max_font * mid_scale
            
            text_size = processor.measure_multiline_text_size(
                text_content, 
                font_name, 
                max_scaled_font,
                available_width_pt - margin_pt * 2
            )
            
            if text_size is None:
                high_scale = mid_scale
                continue
            
            text_width, text_height = text_size
            
            if text_height <= (height_pt - margin_pt * 2) and text_width <= (available_width_pt - margin_pt * 2):
                best_scale = mid_scale
                low_scale = mid_scale
            else:
                high_scale = mid_scale
            
            if high_scale - low_scale < 0.01:
                break
        
        layout['scale_factor'] = best_scale
        print(f'      -> scale={best_scale:.2f} after {iterations} iterations', flush=True)
    
    # Apply
    for layout in shape_layout:
        item = layout['item']
        shape = item['shape']
        original_sizes = item['original_sizes']
        scale_factor = layout.get('scale_factor', 1.0)
        
        print(f'    Applying scale {scale_factor:.2f} to "{shape.name}"...', flush=True)
        for key, orig_size in list(original_sizes.items())[:3]:
            new_size = round(orig_size * scale_factor)
            print(f'      {key}: {orig_size}pt -> {new_size}pt', flush=True)

print('\n=== DONE ===', flush=True)
