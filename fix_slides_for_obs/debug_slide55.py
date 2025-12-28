from pptx import Presentation
import fix_slides_for_obs_processor as processor

# Check slide 55 text in detail
prs = Presentation('Apresentação1Original.pptx')
slide = list(prs.slides)[54]  # 0-indexed

# Get slide dimensions
slide_width = prs.slide_width / 12700  # in points
slide_height = prs.slide_height / 12700  # in points

margin_pt = 10
safety_factor = 0.98  # For single shape
usable_width = (slide_width - 2 * margin_pt) * safety_factor
usable_height = (slide_height - 2 * margin_pt) * safety_factor

print(f'Slide dimensions: {slide_width:.0f}pt x {slide_height:.0f}pt')
print(f'Usable area: {usable_width:.0f}pt x {usable_height:.0f}pt')

for shape in slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = processor.normalize_text_whitespace(shape.text_frame.text)
        
        # Get original font info
        orig_font_size = None
        font_name = 'Arial'
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() and run.font.size:
                    orig_font_size = run.font.size.pt
                    if run.font.name:
                        font_name = run.font.name
                    break
            if orig_font_size:
                break
        
        print(f'\n=== Text Analysis ===')
        print(f'Original font: {orig_font_size}pt')
        print(f'Font name: {font_name}')
        print(f'Text: "{text}"')
        print(f'Text length: {len(text)} chars')
        
        # Test different scale factors
        print('\n=== Scale Factor Tests ===')
        for scale in [1.0, 1.5, 2.0, 2.5, 3.0]:
            test_font = orig_font_size * scale
            text_size = processor.measure_multiline_text_size(text, font_name, test_font, usable_width)
            if text_size:
                fits_width = text_size[0] <= usable_width
                fits_height = text_size[1] <= usable_height
                print(f'Scale {scale:.1f} ({test_font:.0f}pt): {text_size[0]:.0f}x{text_size[1]:.0f}pt '
                      f'- Width {"OK" if fits_width else "OVERFLOW"}, Height {"OK" if fits_height else "OVERFLOW"}')
