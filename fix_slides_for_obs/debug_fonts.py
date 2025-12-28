from pptx import Presentation
from pptx.util import Pt

prs = Presentation('Apresentação1Original.pptx')

print('=== ORIGINAL FONT SIZES ===')
for slide_num, slide in enumerate(prs.slides, 1):
    print(f'\nSlide {slide_num}:')
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f'  Shape: {shape.name}')
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    if run.text.strip():
                        size = run.font.size.pt if run.font.size else 'None'
                        text_preview = run.text[:30].replace('\n', '\\n')
                        print(f'    [{para_idx},{run_idx}] "{text_preview}" -> {size}pt')
