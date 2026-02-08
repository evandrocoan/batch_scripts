"""
Shared functions for PowerPoint slide processing with OBS chroma key support.
"""
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE

try:
    from PIL import ImageFont, ImageDraw, Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

# Common Windows fonts paths
import os
WINDOWS_FONTS_DIR = os.path.join(os.environ.get('WINDIR', 'C:\\Windows'), 'Fonts')

# Default/fallback font
DEFAULT_FONT = 'arial.ttf'


# Placeholder types that should be ignored when checking for meaningful text content
# These are typically auto-generated elements like page numbers, dates, and footers
INSIGNIFICANT_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.HEADER,
}


def is_insignificant_placeholder(shape):
    """
    Check if a shape is an insignificant placeholder (page number, footer, date, header).
    These should be ignored when determining if a slide has meaningful content.
    
    Args:
        shape: A PowerPoint shape object
    
    Returns:
        bool: True if the shape is an insignificant placeholder, False otherwise
    """
    if not shape.is_placeholder:
        return False
    
    try:
        placeholder_format = shape.placeholder_format
        if placeholder_format and placeholder_format.type in INSIGNIFICANT_PLACEHOLDER_TYPES:
            return True
    except Exception:
        pass
    
    return False


def slide_has_visual_elements(slide):
    """
    Check if a slide contains any visual elements like pictures, charts, videos, etc.
    These elements typically have floating positions and text should not be repositioned.
    
    Args:
        slide: A PowerPoint slide object
    
    Returns:
        bool: True if the slide has at least one visual element, False otherwise
    """
    # Shape types that indicate visual/floating elements
    visual_types = {
        MSO_SHAPE_TYPE.PICTURE,        # Images
        MSO_SHAPE_TYPE.CHART,          # Charts/graphs
        MSO_SHAPE_TYPE.MEDIA,          # Videos and audio
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,  # Embedded objects (Excel charts, etc.)
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,    # Linked objects
        MSO_SHAPE_TYPE.TABLE,          # Tables
        MSO_SHAPE_TYPE.GROUP,          # Grouped shapes (often contain visual elements)
        MSO_SHAPE_TYPE.DIAGRAM,        # SmartArt diagrams
        MSO_SHAPE_TYPE.CANVAS,         # Drawing canvas
        MSO_SHAPE_TYPE.INK,            # Ink drawings
        MSO_SHAPE_TYPE.INK_COMMENT,    # Ink comments
        MSO_SHAPE_TYPE.WEB_VIDEO,      # Web videos
        MSO_SHAPE_TYPE.IGX_GRAPHIC,    # SmartArt graphics
    }
    
    for shape in slide.shapes:
        if shape.shape_type in visual_types:
            return True
    return False


import re

def normalize_text_whitespace(text):
    """
    Normalize whitespace in text:
    - Strip leading/trailing whitespace
    - Replace multiple consecutive blank lines with a single blank line
    
    Args:
        text: The text to normalize
    
    Returns:
        str: Normalized text
    """
    if not text:
        return text
    
    # Strip leading/trailing whitespace
    text = text.strip()
    
    # Replace multiple consecutive blank lines (2+ newlines) with a single blank line
    # This pattern matches 2 or more newlines (with optional whitespace between them)
    text = re.sub(r'(\n\s*){2,}', '\n\n', text)
    
    return text


def clean_empty_paragraphs(text_frame):
    """
    Remove empty paragraphs from a text frame, keeping only one blank line between content.
    Also removes trailing empty paragraphs.
    
    Args:
        text_frame: A PowerPoint text frame object
    """
    # Access the XML element directly to remove empty paragraphs
    txBody = text_frame._txBody
    if txBody is None:
        return
    
    # Get all paragraph elements
    p_elements = txBody.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}p')
    
    if not p_elements:
        return
    
    # Track which paragraphs to remove
    paragraphs_to_remove = []
    consecutive_empty = 0
    
    for i, p_elem in enumerate(p_elements):
        # Check if paragraph is empty (no text content)
        text_content = ''.join(p_elem.itertext()).strip()
        
        if not text_content:
            consecutive_empty += 1
            # Remove if: more than one consecutive empty, or it's a trailing empty
            if consecutive_empty > 1:
                paragraphs_to_remove.append(p_elem)
        else:
            consecutive_empty = 0
    
    # Remove trailing empty paragraphs
    for i in range(len(p_elements) - 1, -1, -1):
        p_elem = p_elements[i]
        text_content = ''.join(p_elem.itertext()).strip()
        if not text_content:
            if p_elem not in paragraphs_to_remove:
                paragraphs_to_remove.append(p_elem)
        else:
            break  # Stop at first non-empty paragraph from the end
    
    # Remove the marked paragraphs
    for p_elem in paragraphs_to_remove:
        parent = p_elem.getparent()
        if parent is not None:
            parent.remove(p_elem)


def get_font_path(font_name):
    """
    Get the path to a font file based on font name.
    
    Args:
        font_name: Name of the font (e.g., 'Arial', 'Calibri')
    
    Returns:
        Path to the font file, or default font if not found
    """
    if not font_name:
        font_name = 'Arial'
    
    # Common font name to file mappings
    font_mappings = {
        'arial': 'arial.ttf',
        'arial black': 'ariblk.ttf',
        'calibri': 'calibri.ttf',
        'calibri light': 'calibril.ttf',
        'times new roman': 'times.ttf',
        'verdana': 'verdana.ttf',
        'tahoma': 'tahoma.ttf',
        'trebuchet ms': 'trebuc.ttf',
        'georgia': 'georgia.ttf',
        'comic sans ms': 'comic.ttf',
        'impact': 'impact.ttf',
        'courier new': 'cour.ttf',
        'consolas': 'consola.ttf',
        'segoe ui': 'segoeui.ttf',
    }
    
    font_lower = font_name.lower()
    font_file = font_mappings.get(font_lower, f'{font_lower}.ttf')
    font_path = os.path.join(WINDOWS_FONTS_DIR, font_file)
    
    if os.path.exists(font_path):
        return font_path
    
    # Try default font
    default_path = os.path.join(WINDOWS_FONTS_DIR, DEFAULT_FONT)
    if os.path.exists(default_path):
        return default_path
    
    return None


def measure_text_size(text, font_name, font_size_pt):
    """
    Measure the rendered size of text using Pillow.
    
    Args:
        text: The text to measure
        font_name: Name of the font
        font_size_pt: Font size in points
    
    Returns:
        tuple: (width, height) in pixels, or None if measurement fails
    """
    if not PILLOW_AVAILABLE:
        return None
    
    font_path = get_font_path(font_name)
    if not font_path:
        return None
    
    try:
        font = ImageFont.truetype(font_path, int(font_size_pt))
        # Create a temporary image to get text bbox
        img = Image.new('RGB', (1, 1))
        draw = ImageDraw.Draw(img)
        bbox = draw.textbbox((0, 0), text, font=font)
        width = bbox[2] - bbox[0]
        height = bbox[3] - bbox[1]
        return (width, height)
    except Exception:
        return None


def measure_multiline_text_size(text, font_name, font_size_pt, max_width=None):
    """
    Measure the rendered size of potentially multi-line text.
    
    Args:
        text: The text to measure (may contain newlines)
        font_name: Name of the font
        font_size_pt: Font size in points
        max_width: Maximum width for word wrapping (pixels), None for no wrapping
    
    Returns:
        tuple: (width, height) in pixels, or None if measurement fails
    """
    if not PILLOW_AVAILABLE:
        return None
    
    font_path = get_font_path(font_name)
    if not font_path:
        return None
    
    try:
        font = ImageFont.truetype(font_path, int(font_size_pt))
        img = Image.new('RGB', (1, 1))
        draw = ImageDraw.Draw(img)
        
        lines = text.split('\n')
        total_height = 0
        max_line_width = 0
        
        # PowerPoint uses larger line spacing than Pillow measures
        # Use line_spacing_factor to better match PowerPoint rendering
        # 0.2 = too aggressive (text overflows), 0.5 = too conservative (too much space)
        # 0.35 = still caused overflow on some slides, 0.40 = safer
        line_spacing_factor = 0.40
        
        for line in lines:
            if not line.strip():
                # Empty line - add line height
                total_height += font_size_pt * 1.35
                continue
            
            # If max_width specified, wrap the line
            if max_width:
                words = line.split()
                current_line = ""
                for word in words:
                    test_line = f"{current_line} {word}".strip()
                    bbox = draw.textbbox((0, 0), test_line, font=font)
                    line_width = bbox[2] - bbox[0]
                    
                    if line_width <= max_width:
                        current_line = test_line
                    else:
                        if current_line:
                            bbox = draw.textbbox((0, 0), current_line, font=font)
                            max_line_width = max(max_line_width, bbox[2] - bbox[0])
                            total_height += bbox[3] - bbox[1] + font_size_pt * line_spacing_factor
                        current_line = word
                
                if current_line:
                    bbox = draw.textbbox((0, 0), current_line, font=font)
                    max_line_width = max(max_line_width, bbox[2] - bbox[0])
                    total_height += bbox[3] - bbox[1] + font_size_pt * line_spacing_factor
            else:
                bbox = draw.textbbox((0, 0), line, font=font)
                max_line_width = max(max_line_width, bbox[2] - bbox[0])
                total_height += bbox[3] - bbox[1] + font_size_pt * line_spacing_factor
        
        return (max_line_width, total_height)
    except Exception:
        return None


def check_text_overflow(shape, slide_width, slide_height):
    """
    Check if a shape's text overflows the slide boundaries.
    
    Args:
        shape: PowerPoint shape object
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    
    Returns:
        dict: {
            'overflows': bool,
            'shape_right': int (EMUs),
            'shape_bottom': int (EMUs),
            'overflow_right': int (EMUs, negative if within bounds),
            'overflow_bottom': int (EMUs, negative if within bounds)
        }
    """
    shape_left = shape.left if shape.left else 0
    shape_top = shape.top if shape.top else 0
    shape_width = shape.width if shape.width else 0
    shape_height = shape.height if shape.height else 0
    
    shape_right = shape_left + shape_width
    shape_bottom = shape_top + shape_height
    
    overflow_right = shape_right - slide_width
    overflow_bottom = shape_bottom - slide_height
    overflow_left = -shape_left if shape_left < 0 else 0
    overflow_top = -shape_top if shape_top < 0 else 0
    
    overflows = (overflow_right > 0 or overflow_bottom > 0 or 
                 overflow_left > 0 or overflow_top > 0)
    
    return {
        'overflows': overflows,
        'shape_right': shape_right,
        'shape_bottom': shape_bottom,
        'overflow_right': overflow_right,
        'overflow_bottom': overflow_bottom,
        'overflow_left': overflow_left,
        'overflow_top': overflow_top
    }


def get_shape_font_info(shape):
    """
    Get font information from a shape's text.
    
    Args:
        shape: PowerPoint shape with text
    
    Returns:
        dict: {'font_name': str, 'font_size': Pt or None, 'text': str}
    """
    font_name = 'Arial'
    font_size = None
    full_text = ""
    
    if not shape.has_text_frame:
        return {'font_name': font_name, 'font_size': font_size, 'text': full_text}
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            full_text += run.text
            if run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size = run.font.size
        full_text += '\n'
    
    return {'font_name': font_name, 'font_size': font_size, 'text': full_text.strip()}


def calculate_max_font_size(shape, slide_width, slide_height, margin_pt=10, min_size=8, max_size=200):
    """
    Calculate the maximum font size that fits text within the shape without overflow.
    Uses binary search for efficiency.
    
    Args:
        shape: PowerPoint shape with text
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
        margin_pt: Margin in points to leave around text
        min_size: Minimum font size to try
        max_size: Maximum font size to try
    
    Returns:
        int: Maximum font size in points that fits, or None if calculation fails
    """
    if not PILLOW_AVAILABLE:
        return None
    
    font_info = get_shape_font_info(shape)
    text = font_info['text']
    font_name = font_info['font_name']
    
    if not text.strip():
        return None
    
    # Get shape dimensions in points (1 point = 12700 EMUs)
    shape_width_pt = (shape.width / 12700) if shape.width else 0
    shape_height_pt = (shape.height / 12700) if shape.height else 0
    
    # Account for margins
    available_width = shape_width_pt - (margin_pt * 2)
    available_height = shape_height_pt - (margin_pt * 2)
    
    if available_width <= 0 or available_height <= 0:
        return min_size
    
    # Binary search for maximum font size
    low = min_size
    high = max_size
    best_size = min_size
    
    while low <= high:
        mid = (low + high) // 2
        
        text_size = measure_multiline_text_size(text, font_name, mid, available_width)
        
        if text_size is None:
            return None
        
        text_width, text_height = text_size
        
        # Check if text fits
        if text_width <= available_width and text_height <= available_height:
            best_size = mid
            low = mid + 1  # Try larger
        else:
            high = mid - 1  # Try smaller
    
    return best_size


def auto_fit_text_to_shape(shape, slide_width, slide_height, margin_pt=10):
    """
    Automatically adjust font size to maximize text size while fitting in shape.
    
    Args:
        shape: PowerPoint shape with text
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
        margin_pt: Margin in points to leave around text
    
    Returns:
        int: New font size applied, or None if no change made
    """
    max_font_size = calculate_max_font_size(shape, slide_width, slide_height, margin_pt)
    
    if max_font_size is None:
        return None
    
    # Apply the new font size to all runs
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(max_font_size)
    
    return max_font_size


def check_and_report_overflow(prs):
    """
    Check all slides for text overflow and return a report.
    
    Args:
        prs: PowerPoint Presentation object
    
    Returns:
        list: List of overflow reports [{slide_num, shape_name, overflow_info}]
    """
    overflow_report = []
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            if not shape.text_frame.text.strip():
                continue
            
            overflow_info = check_text_overflow(shape, slide_width, slide_height)
            
            if overflow_info['overflows']:
                overflow_report.append({
                    'slide_num': slide_num,
                    'shape_name': shape.name if hasattr(shape, 'name') else 'Unknown',
                    'overflow_info': overflow_info
                })
    
    return overflow_report


def auto_fit_all_text(prs, margin_pt=10):
    """
    Automatically fit all text in the presentation to maximum size without overflow.
    
    Args:
        prs: PowerPoint Presentation object
        margin_pt: Margin in points to leave around text
    
    Returns:
        list: List of changes made [{slide_num, shape_name, old_size, new_size}]
    """
    if not PILLOW_AVAILABLE:
        raise ImportError("Pillow is required for auto-fit. Install with: pip install Pillow")
    
    changes = []
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            if not shape.text_frame.text.strip():
                continue
            
            # Get current font size
            font_info = get_shape_font_info(shape)
            old_size = font_info['font_size']
            old_size_pt = old_size.pt if old_size else None
            
            # Calculate and apply new size
            new_size = auto_fit_text_to_shape(shape, slide_width, slide_height, margin_pt)
            
            if new_size is not None:
                changes.append({
                    'slide_num': slide_num,
                    'shape_name': shape.name if hasattr(shape, 'name') else 'Unknown',
                    'old_size': old_size_pt,
                    'new_size': new_size
                })
    
    return changes


def reposition_and_resize_text_boxes(prs, margin_percent=0.05, spacing_pt=10):
    """
    Reposition all text boxes on each slide to be stacked vertically and fill the page.
    Text boxes are sized proportionally based on their text content length.
    
    Args:
        prs: PowerPoint Presentation object
        margin_percent: Margin as percentage of slide dimensions (0.05 = 5%)
        spacing_pt: Spacing between text boxes in points
    
    Returns:
        int: Number of slides processed
    """
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Calculate margins in EMUs
    margin_x = int(slide_width * margin_percent)
    margin_y = int(slide_height * margin_percent)
    spacing_emu = int(spacing_pt * 12700)  # Convert points to EMUs
    
    # Available area
    available_width = slide_width - (2 * margin_x)
    available_height = slide_height - (2 * margin_y)
    
    slides_processed = 0
    
    for slide in prs.slides:
        # Skip slides with visual elements (images, charts, videos, etc. shouldn't have text repositioned)
        if slide_has_visual_elements(slide):
            continue
        
        # Collect all text shapes with content and their text lengths
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                # Clean empty paragraphs from the text frame first
                clean_empty_paragraphs(shape.text_frame)
                
                text = normalize_text_whitespace(shape.text_frame.text)
                # Count lines and characters to estimate needed height
                lines = text.count('\n') + 1
                char_count = len(text)
                # Weight: combination of line count and character count
                # More lines = needs more height, more chars per line also needs height
                weight = max(1, lines + (char_count / 50))  # Roughly 50 chars per line estimate
                text_shapes.append({'shape': shape, 'weight': weight, 'text': text})
        
        if not text_shapes:
            continue
        
        slides_processed += 1
        num_shapes = len(text_shapes)
        
        # Calculate total weight
        total_weight = sum(item['weight'] for item in text_shapes)
        
        # Calculate total spacing needed
        total_spacing = spacing_emu * (num_shapes - 1) if num_shapes > 1 else 0
        
        # Height available for text boxes
        height_for_boxes = available_height - total_spacing
        
        # Reposition each text box with proportional height
        current_y = margin_y
        
        for item in text_shapes:
            shape = item['shape']
            weight = item['weight']
            
            # Calculate proportional height based on weight
            height_ratio = weight / total_weight
            height_per_box = int(height_for_boxes * height_ratio)
            
            # Ensure minimum height
            min_height = int(available_height * 0.1)  # At least 10% of available height
            height_per_box = max(height_per_box, min_height)
            
            try:
                # Set position and size
                shape.left = margin_x
                shape.top = current_y
                shape.width = available_width
                shape.height = height_per_box
                
                # Enable "Shrink text on overflow" to prevent text from going outside the box
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # Center align text in the shape
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                
                # Move to next position
                current_y += height_per_box + spacing_emu
            except Exception:
                pass  # Some shapes may not allow repositioning
    
    return slides_processed


def reposition_and_maximize_font(prs, margin_percent=0.05, spacing_pt=10):
    """
    Reposition text boxes to fill the page and maximize font size.
    
    Algorithm:
    1. First, resize text boxes based on content (lines) - WITHOUT changing fonts
       e.g., if shape1 has 1 line and shape2 has 2 lines, shape1 gets 1/3, shape2 gets 2/3
    2. Then, calculate a single scale factor that fits ALL text in their boxes
    3. Apply that scale factor to all fonts (preserving original proportions)
    
    Args:
        prs: PowerPoint Presentation object
        margin_percent: Margin as percentage of slide dimensions (0.05 = 5%)
        spacing_pt: Spacing between text boxes in points
    
    Returns:
        dict: {'slides_processed': int, 'font_changes': list}
    """
    if not PILLOW_AVAILABLE:
        raise ImportError("Pillow is required for font maximization. Install with: pip install Pillow")
    
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Calculate margins in EMUs
    margin_x = int(slide_width * margin_percent)
    margin_y = int(slide_height * margin_percent)
    spacing_emu = int(spacing_pt * 12700)
    
    # Available area
    available_width = slide_width - (2 * margin_x)
    available_height = slide_height - (2 * margin_y)
    available_width_pt = available_width / 12700
    available_height_pt = available_height / 12700
    
    slides_processed = 0
    font_changes = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Skip slides with visual elements (images, charts, videos, etc. shouldn't have text repositioned)
        if slide_has_visual_elements(slide):
            continue
        
        # STEP 1: Collect all shapes and their original font sizes
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                # Clean empty paragraphs from the text frame first
                clean_empty_paragraphs(shape.text_frame)
                
                text = normalize_text_whitespace(shape.text_frame.text)
                
                # Store original position for sorting
                original_top = shape.top
                
                # Store original sizes indexed by position
                original_sizes = {}  # (para_idx, run_idx) -> original_size
                max_font_in_shape = 0
                font_name = 'Arial'
                
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            if run.font.size:
                                font_size = run.font.size.pt
                            else:
                                font_size = 12  # Default
                            
                            original_sizes[(para_idx, run_idx)] = font_size
                            max_font_in_shape = max(max_font_in_shape, font_size)
                            
                            if run.font.name:
                                font_name = run.font.name
                
                if max_font_in_shape == 0:
                    max_font_in_shape = 12
                
                # Count lines at ORIGINAL font size for weight calculation
                # This determines box height distribution
                text_size = measure_multiline_text_size(
                    text, font_name, max_font_in_shape, available_width_pt - 20
                )
                if text_size:
                    # Weight = measured height at original font size
                    weight = text_size[1]
                else:
                    # Fallback: simple line count * font size
                    lines = text.count('\n') + 1
                    weight = max_font_in_shape * lines
                
                text_shapes.append({
                    'shape': shape, 
                    'weight': weight,
                    'text': text,
                    'original_sizes': original_sizes,
                    'max_font': max_font_in_shape,
                    'font_name': font_name,
                    'original_top': original_top  # Store for sorting
                })
        
        if not text_shapes:
            continue
        
        # Sort shapes by their original vertical position (top to bottom)
        # This preserves the visual order from the original slide
        text_shapes.sort(key=lambda x: x['original_top'])
        
        slides_processed += 1
        num_shapes = len(text_shapes)
        
        # STEP 2: Calculate text box heights based on content weight
        total_weight = sum(item['weight'] for item in text_shapes)
        total_spacing = spacing_emu * (num_shapes - 1) if num_shapes > 1 else 0
        height_for_boxes = available_height - total_spacing
        
        current_y = margin_y
        shape_layout = []
        
        for item in text_shapes:
            height_ratio = item['weight'] / total_weight
            height_per_box = int(height_for_boxes * height_ratio)
            min_height = int(available_height * 0.1)
            height_per_box = max(height_per_box, min_height)
            
            shape_layout.append({
                'item': item,
                'height': height_per_box,
                'height_pt': height_per_box / 12700,
                'y': current_y
            })
            current_y += height_per_box + spacing_emu
        
        # STEP 3: Find the maximum scale factor that fits ALL shapes
        # Binary search for a single global scale factor
        # IMPORTANT: Never shrink below original size (scale >= 1.0)
        margin_pt = 10  # Margin for text box padding
        # Use smaller safety factor for single shapes, larger for multiple
        safety_factor = 0.98 if num_shapes == 1 else 0.95
        low_scale = 1.0  # Never shrink fonts below original
        high_scale = 20.0
        best_scale = 1.0  # Default: keep original size

        # Calculate usable dimensions once (same for wrap and check)
        usable_width = (available_width_pt - margin_pt * 2) * safety_factor
        
        for _ in range(25):  # Binary search iterations
            mid_scale = (low_scale + high_scale) / 2
            all_fit = True
            
            # Check if ALL shapes fit at this scale
            for layout in shape_layout:
                item = layout['item']
                height_pt = layout['height_pt']
                max_font = item['max_font']
                font_name = item['font_name']
                text_content = item['text']
                
                # Scale the largest font in this shape
                max_scaled_font = max_font * mid_scale
                
                # Measure text at this scaled font size
                # Use usable_width for wrapping to ensure wrapped text fits in usable area
                text_size = measure_multiline_text_size(
                    text_content, 
                    font_name, 
                    max_scaled_font,
                    usable_width
                )
                
                if text_size is None:
                    all_fit = False
                    break
                
                text_width, text_height = text_size
                
                # Check if it fits in this shape's allocated box
                # Use safety_factor to account for PowerPoint's actual text rendering
                usable_height = (height_pt - margin_pt * 2) * safety_factor
                if text_height > usable_height or text_width > usable_width:
                    all_fit = False
                    break
            
            if all_fit:
                best_scale = mid_scale
                low_scale = mid_scale
            else:
                high_scale = mid_scale
            
            if high_scale - low_scale < 0.01:
                break
        
        # Print debug info for slide 57 specifically, or slides that didn't scale much
        if slide_num == 57 or best_scale < 1.5:
            print(f"DEBUG Slide {slide_num}: best_scale={best_scale:.2f}, num_shapes={num_shapes}")
            for layout in shape_layout:
                item = layout['item']
                print(f"  - height_pt={layout['height_pt']:.1f}, max_font={item['max_font']}, text='{item['text'][:50]}...'")
                # Test measurement at scale 2.0
                test_size = measure_multiline_text_size(
                    item['text'], item['font_name'], item['max_font'] * 2.0,
                    usable_width
                )
                if test_size:
                    usable_h = (layout['height_pt'] - margin_pt * 2) * safety_factor
                    print(f"    at scale 2.0: text_height={test_size[1]:.1f}, usable_height={usable_h:.1f}")
        
        # STEP 4: Apply positions, sizes, and scaled fonts
        for layout in shape_layout:
            item = layout['item']
            shape = item['shape']
            original_sizes = item['original_sizes']
            
            try:
                # Set text box position and size
                shape.left = margin_x
                shape.top = layout['y']
                shape.width = available_width
                shape.height = layout['height']
                
                # Enable "Shrink text on overflow" to prevent text from going outside the box
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # Apply scaled font sizes (same scale for all, preserving ratios)
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            key = (para_idx, run_idx)
                            if key in original_sizes:
                                original_size = original_sizes[key]
                                new_size = round(original_size * best_scale)
                                new_size = max(8, min(200, new_size))
                                run.font.size = Pt(new_size)
                                
                                font_changes.append({
                                    'slide_num': slide_num,
                                    'shape_name': shape.name if hasattr(shape, 'name') else 'Unknown',
                                    'old_size': original_size,
                                    'new_size': new_size,
                                    'scale': best_scale
                                })
            except Exception:
                pass
    
    return {'slides_processed': slides_processed, 'font_changes': font_changes}


from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


def auto_fit_and_center_text_frame(text_frame, slide_width, slide_height, margin=0.05):
    """
    Auto-fit text by reducing font size and center the text frame to prevent overflow.
    
    Args:
        text_frame: The text frame to process
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
        margin: Margin as percentage of slide dimensions (0.05 = 5%)
    """
    if not text_frame.text.strip():
        return
    
    # Calculate available space with margins
    max_width = slide_width * (1 - 2 * margin)
    max_height = slide_height * (1 - 2 * margin)
    
    # Get the shape containing this text frame
    shape = text_frame._element.getparent().getparent()
    
    # Find minimum and maximum font sizes across all runs
    min_font_size = 8  # Minimum readable size
    max_font_size = 72  # Starting maximum
    
    # Get initial font sizes
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                font_pt = run.font.size.pt
                if font_pt > max_font_size:
                    max_font_size = font_pt
    
    # Binary search for optimal font size
    current_size = max_font_size
    optimal_size = min_font_size
    
    for attempt in range(10):  # Limit iterations
        # Set all text to current size
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    run.font.size = Pt(current_size)
        
        # Check if text fits (simplified check using character estimation)
        total_chars = len(text_frame.text)
        estimated_char_width = current_size * 0.6  # Rough estimation
        estimated_width = total_chars * estimated_char_width * 12700  # Convert to EMUs
        
        lines = text_frame.text.count('\n') + 1
        estimated_height = lines * current_size * 1.2 * 12700  # Line height with spacing
        
        if estimated_width <= max_width and estimated_height <= max_height:
            optimal_size = current_size
            if current_size == max_font_size:
                break  # Already at maximum
            current_size = min(max_font_size, current_size + (max_font_size - current_size) / 2)
        else:
            max_font_size = current_size
            current_size = (current_size + min_font_size) / 2
        
        if abs(max_font_size - min_font_size) < 1:
            break
    
    # Apply optimal size
    for paragraph in text_frame.paragraphs:
        # Center align paragraphs
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            if run.text.strip():
                run.font.size = Pt(optimal_size)
    
    # Center the text frame itself
    try:
        # Calculate centered position
        shape_width = max_width
        shape_height = max_height
        
        # Position shape in center of slide
        shape.left = int((slide_width - shape_width) / 2)
        shape.top = int((slide_height - shape_height) / 2)
        shape.width = int(shape_width)
        shape.height = int(shape_height)
    except:
        pass  # Some shapes may not allow positioning


def apply_solid_glow_to_run(run, color_hex, size_pt):
    """
    Applies a SOLID glow effect (0% transparency) to create a highlighter background effect.
    Uses the effectLst structure that PowerPoint expects.
    
    Args:
        run: The text run to apply the glow to
        color_hex: Hex color code (with or without #)
        size_pt: Glow size in points
    """
    # Strip '#' from color if present (for VS Code color preview support)
    color_hex = color_hex.lstrip('#')
    
    # Get the Run Properties (rPr) element or create it
    rPr = run._r.get_or_add_rPr()
    
    # Calculate radius in EMUs (1 point = 12700 EMUs)
    radius_emu = int(size_pt * 12700)
    
    # Create the XML for the effect list with glow
    # Using alpha val="100000" for 100% opacity (0% transparency)
    effectlst_xml = f'''<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:glow rad="{radius_emu}">
            <a:srgbClr val="{color_hex}">
                <a:alpha val="100000"/>
            </a:srgbClr>
        </a:glow>
    </a:effectLst>'''
    
    # Remove any existing effects more robustly
    # Check for all possible effect-related elements regardless of namespace
    # This handles different DrawingML versions and effect types
    elements_to_remove = []
    for child in rPr:
        tag = child.tag
        # Check if tag contains effect-related element names (namespace-agnostic)
        if any(effect_name in tag for effect_name in ['effectLst', 'glow', 'outerShdw', 'innerShdw', 'reflection', 'softEdge', 'effectDag']):
            elements_to_remove.append(child)
    
    # Remove all found effect elements
    for element in elements_to_remove:
        rPr.remove(element)
    
    # Parse the new XML and inject it into the run properties
    # effectLst should come after solidFill but before font typefaces
    effectlst_element = parse_xml(effectlst_xml)
    
    # Insert after solidFill if it exists, otherwise at the beginning
    solid_fill = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    if solid_fill is not None:
        # Insert after solidFill
        insert_index = list(rPr).index(solid_fill) + 1
        rPr.insert(insert_index, effectlst_element)
    else:
        rPr.insert(0, effectlst_element)


def reset_master_slides(prs):
    """
    Reset all master slides and layouts to default formatting by removing effects and backgrounds.
    This removes glow effects, shadows, and other visual effects from master slide text.
    
    Args:
        prs: PowerPoint Presentation object
    
    Returns:
        tuple: (number of master slides processed, number of layouts processed)
    """
    masters_count = 0
    layouts_count = 0
    
    # Process each slide master
    for slide_master in prs.slide_masters:
        masters_count += 1
        
        # Reset master slide background to default (white)
        try:
            background = slide_master.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
        except:
            pass  # Some masters may not allow background modification
        
        # Remove effects from all text in master slide
        for shape in slide_master.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            # Remove all effects from run
                            try:
                                rPr = run._r.get_or_add_rPr()
                                elements_to_remove = []
                                for child in rPr:
                                    tag = child.tag
                                    if any(effect_name in tag for effect_name in ['effectLst', 'glow', 'outerShdw', 'innerShdw', 'reflection', 'softEdge', 'effectDag']):
                                        elements_to_remove.append(child)
                                for element in elements_to_remove:
                                    rPr.remove(element)
                            except:
                                pass
        
        # Process each layout in the master
        for layout in slide_master.slide_layouts:
            layouts_count += 1
            
            # Reset layout background
            try:
                background = layout.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
            except:
                pass
            
            # Remove effects from all text in layout
            for shape in layout.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                try:
                                    rPr = run._r.get_or_add_rPr()
                                    elements_to_remove = []
                                    for child in rPr:
                                        tag = child.tag
                                        if any(effect_name in tag for effect_name in ['effectLst', 'glow', 'outerShdw', 'innerShdw', 'reflection', 'softEdge', 'effectDag']):
                                            elements_to_remove.append(child)
                                    for element in elements_to_remove:
                                        rPr.remove(element)
                                except:
                                    pass
    
    return masters_count, layouts_count


def process_presentation(prs, glow_color, glow_size, text_color):
    """
    Process a PowerPoint presentation to add glow effects and set backgrounds.
    
    Args:
        prs: PowerPoint Presentation object
        glow_color: Hex color for glow (with or without #)
        glow_size: Glow size in points
        text_color: Hex color for text (with or without #)
    
    Returns:
        int: Number of text shapes processed
    """
    # Strip '#' from text color if present
    text_color = text_color.lstrip('#')
    
    count = 0
    
    # Iterate through every slide
    for slide in prs.slides:
        # Check if slide has any meaningful text (ignore page numbers, footers, etc.)
        has_text = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                # Skip insignificant placeholders like page numbers, footers, dates
                if is_insignificant_placeholder(shape):
                    continue
                has_text = True
                break
        
        # Set slide background: white if has text, black if no text
        background = slide.background
        fill = background.fill
        fill.solid()
        if has_text:
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        else:
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
            continue  # Skip processing shapes on empty slides
        
        # Iterate through every shape on the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if not shape.has_text_frame:
                continue
            
            # Process all text regardless of whether it's from master or manual text box
            text_processed = False
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():  # Only process non-empty runs
                        # Apply the SOLID glow (0% transparency) to create highlighter effect
                        apply_solid_glow_to_run(run, glow_color, glow_size)
                        
                        # Set text color
                        run.font.color.rgb = RGBColor(
                            int(text_color[0:2], 16),
                            int(text_color[2:4], 16),
                            int(text_color[4:6], 16)
                        )
                        text_processed = True
            
            if text_processed:
                count += 1
    
    return count
