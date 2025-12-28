from pptx import Presentation
import fix_slides_for_obs_processor as processor
from PIL import ImageFont, ImageDraw, Image

# Check slide 55 text line by line
prs = Presentation('Apresentação1Original.pptx')
slide = list(prs.slides)[54]  # 0-indexed

for shape in slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = processor.normalize_text_whitespace(shape.text_frame.text)
        
        # Get original font info
        orig_font_size = 42.86
        font_name = 'Arial'
        
        font_path = processor.get_font_path(font_name)
        
        print(f'Text (total {len(text)} chars):')
        print(f'"{text}"')
        print()
        
        # Measure each line at different font sizes
        lines = text.split('\n')
        print(f'Number of lines: {len(lines)}')
        print()
        
        for scale in [1.0, 1.5, 2.0]:
            test_font_pt = orig_font_size * scale
            font = ImageFont.truetype(font_path, int(test_font_pt))
            img = Image.new('RGB', (1, 1))
            draw = ImageDraw.Draw(img)
            
            print(f'=== Scale {scale:.1f} ({test_font_pt:.0f}pt) ===')
            for i, line in enumerate(lines):
                bbox = draw.textbbox((0, 0), line, font=font)
                width = bbox[2] - bbox[0]
                print(f'  Line {i+1} ({len(line)} chars): {width:.0f}pt wide - "{line[:60]}..."')
            
            # Get total width (considering word wrap to 827pt limit)
            usable_width = 827
            text_size = processor.measure_multiline_text_size(
                text, font_name, test_font_pt, usable_width
            )
            if text_size:
                print(f'  -> After wrapping to {usable_width}pt: {text_size[0]:.0f}x{text_size[1]:.0f}pt')
            print()
