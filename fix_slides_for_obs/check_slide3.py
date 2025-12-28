from pptx import Presentation

# Check slide 3 (ORAÇÃO VOCACIONAL)
prs_orig = Presentation('Apresentação1Original.pptx')
prs_proc = Presentation('test_output.pptx')

slide_orig = list(prs_orig.slides)[2]  # 0-indexed, slide 3
slide_proc = list(prs_proc.slides)[2]

print('=== SLIDE 3 ORIGINAL ===')
shapes_orig = []
for shape in slide_orig.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text.strip()[:50]
        shapes_orig.append((shape.top, shape.name, text))

# Sort by top position
shapes_orig.sort(key=lambda x: x[0])
for top, name, text in shapes_orig:
    print(f'  top={top/914400:.2f}" - {name}: "{text}"')

print('\n=== SLIDE 3 AFTER PROCESSING ===')
shapes_proc = []
for shape in slide_proc.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text.strip()[:50]
        shapes_proc.append((shape.top, shape.name, text))

# Sort by top position
shapes_proc.sort(key=lambda x: x[0])
for top, name, text in shapes_proc:
    print(f'  top={top/914400:.2f}" - {name}: "{text}"')
