from pptx import Presentation
import fix_slides_for_obs_processor as processor

# Debug slide 57 scaling
prs_orig = Presentation('Apresentação1Original.pptx')

slide_orig = list(prs_orig.slides)[56]  # 0-indexed

slide_width = prs_orig.slide_width
slide_height = prs_orig.slide_height

margin_percent = 0.05
margin_x = int(slide_width * margin_percent)
margin_y = int(slide_height * margin_percent)

available_width = slide_width - (2 * margin_x)
available_height = slide_height - (2 * margin_y)
available_width_pt = available_width / 12700
available_height_pt = available_height / 12700

print(f'Slide dimensions: {slide_width/914400:.2f}" x {slide_height/914400:.2f}"')
print(f'Available area: {available_width_pt:.0f}pt x {available_height_pt:.0f}pt')

for shape in slide_orig.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text.strip()
        
        # Get original font info
        max_font = 0
        font_name = 'Arial'
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    if run.font.size:
                        max_font = max(max_font, run.font.size.pt)
                    if run.font.name:
                        font_name = run.font.name
        
        print(f'\nShape: {shape.name}')
        print(f'Original font: {max_font}pt, font_name: {font_name}')
        print(f'Text: "{text}"')
        
        # Try different scales and see what fits
        margin_pt = 10
        safety_factor = 0.98
        usable_height = (available_height_pt - margin_pt * 2) * safety_factor
        usable_width = (available_width_pt - margin_pt * 2) * safety_factor
        
        print(f'\nUsable area: {usable_width:.0f}pt x {usable_height:.0f}pt')
        
        print('\nScale tests:')
        for scale in [1.0, 1.5, 2.0, 2.5, 3.0, 3.5, 4.0]:
            scaled_font = max_font * scale
            text_size = processor.measure_multiline_text_size(
                text, font_name, scaled_font, available_width_pt - margin_pt * 2
            )
            if text_size:
                fits_h = "OK" if text_size[1] <= usable_height else "OVERFLOW"
                fits_w = "OK" if text_size[0] <= usable_width else "OVERFLOW"
                print(f'  Scale {scale:.1f}: font={scaled_font:.0f}pt, size={text_size[0]:.0f}x{text_size[1]:.0f}pt, height:{fits_h}, width:{fits_w}')
