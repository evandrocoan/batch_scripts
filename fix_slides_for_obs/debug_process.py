from pptx import Presentation
from pptx.util import Pt
import fix_slides_for_obs_processor as processor

# Make a copy and process
import shutil
shutil.copy('Apresentação1Original.pptx', 'test_output.pptx')

prs = Presentation('test_output.pptx')

print('=== PROCESSING WITH REPOSITION ===')
result = processor.reposition_and_maximize_font(prs, margin_percent=0.05, spacing_pt=10)

print(f'\nSlides processed: {result["slides_processed"]}')
print(f'\nFont changes (first 30):')
for change in result['font_changes'][:30]:
    print(f"  Slide {change['slide_num']}: {change['old_size']}pt -> {change['new_size']}pt (scale: {change.get('scale', 'N/A'):.2f})")

# Save and re-read to show final state
prs.save('test_output.pptx')

print('\n\n=== FINAL FONT SIZES (first 10 slides) ===')
prs2 = Presentation('test_output.pptx')
for slide_num, slide in enumerate(prs2.slides, 1):
    if slide_num > 10:
        break
    has_content = False
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            if not has_content:
                print(f'\nSlide {slide_num}:')
                has_content = True
            print(f'  Shape: {shape.name}')
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    if run.text.strip():
                        size = run.font.size.pt if run.font.size else 'None'
                        text_preview = run.text[:30].replace('\n', '\\n')
                        print(f'    [{para_idx},{run_idx}] "{text_preview}" -> {size}pt')
