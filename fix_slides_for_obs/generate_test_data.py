"""
Generate expected test data for all slides.
This script analyzes both original and processed presentations
to capture the expected state for each slide.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import fix_slides_for_obs_processor as processor

def slide_has_visual_elements(slide):
    """Check if slide has images, charts, videos, or other visual elements."""
    visual_types = {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.CHART,
        MSO_SHAPE_TYPE.MEDIA,
        MSO_SHAPE_TYPE.TABLE,
        MSO_SHAPE_TYPE.GROUP,
        MSO_SHAPE_TYPE.DIAGRAM,
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
    }
    for shape in slide.shapes:
        if shape.shape_type in visual_types:
            return True
        if hasattr(shape, 'has_chart') and shape.has_chart:
            return True
        if hasattr(shape, 'has_table') and shape.has_table:
            return True
    return False

def get_slide_text_info(slide):
    """Extract text shape information from a slide."""
    shapes_info = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            shape_info = {
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'text': shape.text_frame.text.strip()[:100],  # First 100 chars
                'fonts': []
            }
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() and run.font.size:
                        shape_info['fonts'].append({
                            'size_pt': run.font.size.pt,
                            'name': run.font.name or 'Arial'
                        })
            shapes_info.append(shape_info)
    return shapes_info

def analyze_all_slides():
    """Analyze all slides and generate test expectations."""
    prs_orig = Presentation('Apresentação1Original.pptx')
    prs_proc = Presentation('test_output.pptx')
    
    slides_orig = list(prs_orig.slides)
    slides_proc = list(prs_proc.slides)
    
    test_data = []
    
    for i, (slide_orig, slide_proc) in enumerate(zip(slides_orig, slides_proc)):
        slide_num = i + 1
        
        has_visual = slide_has_visual_elements(slide_orig)
        orig_info = get_slide_text_info(slide_orig)
        proc_info = get_slide_text_info(slide_proc)
        
        # Determine if slide should be transformed
        should_transform = not has_visual and len(orig_info) > 0
        
        # Check if it was actually transformed
        was_transformed = False
        if orig_info and proc_info:
            # Compare positions - if they changed significantly, it was transformed
            if orig_info[0]['left'] != proc_info[0]['left'] or orig_info[0]['top'] != proc_info[0]['top']:
                was_transformed = True
            # Also check font size changes
            if orig_info[0]['fonts'] and proc_info[0]['fonts']:
                if orig_info[0]['fonts'][0]['size_pt'] != proc_info[0]['fonts'][0]['size_pt']:
                    was_transformed = True
        
        slide_data = {
            'slide_num': slide_num,
            'has_visual_elements': has_visual,
            'has_text': len(orig_info) > 0,
            'should_transform': should_transform,
            'was_transformed': was_transformed,
            'num_text_shapes_orig': len(orig_info),
            'num_text_shapes_proc': len(proc_info),
        }
        
        # Store original font sizes
        if orig_info and orig_info[0]['fonts']:
            slide_data['orig_font_pt'] = orig_info[0]['fonts'][0]['size_pt']
        
        # Store processed font sizes
        if proc_info and proc_info[0]['fonts']:
            slide_data['proc_font_pt'] = proc_info[0]['fonts'][0]['size_pt']
        
        # Store processed position (for transformed slides)
        if was_transformed and proc_info:
            slide_data['proc_left'] = proc_info[0]['left']
            slide_data['proc_top'] = proc_info[0]['top']
            slide_data['proc_width'] = proc_info[0]['width']
            slide_data['proc_height'] = proc_info[0]['height']
        
        test_data.append(slide_data)
        
        # Print summary
        status = "TRANSFORMED" if was_transformed else ("SKIPPED (visual)" if has_visual else "UNCHANGED")
        font_info = ""
        if 'orig_font_pt' in slide_data and 'proc_font_pt' in slide_data:
            font_info = f" ({slide_data['orig_font_pt']:.1f}pt -> {slide_data['proc_font_pt']:.1f}pt)"
        print(f"Slide {slide_num:2d}: {status}{font_info}")
    
    return test_data

if __name__ == '__main__':
    print("Analyzing all slides...\n")
    test_data = analyze_all_slides()
    
    # Save to JSON for test file generation
    with open('test_expected_data.json', 'w', encoding='utf-8') as f:
        json.dump(test_data, f, indent=2, ensure_ascii=False)
    
    print(f"\nSaved {len(test_data)} slide expectations to test_expected_data.json")
    
    # Summary
    transformed = sum(1 for s in test_data if s['was_transformed'])
    skipped_visual = sum(1 for s in test_data if s['has_visual_elements'])
    unchanged = sum(1 for s in test_data if not s['was_transformed'] and not s['has_visual_elements'])
    
    print(f"\nSummary:")
    print(f"  Transformed: {transformed}")
    print(f"  Skipped (visual elements): {skipped_visual}")
    print(f"  Unchanged: {unchanged}")
