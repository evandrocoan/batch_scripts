from pptx import Presentation

# Original
prs = Presentation('Apresentação1Original.pptx')
slide = list(prs.slides)[16]  # 0-indexed, so slide 17

print('=== SLIDE 17 ORIGINAL ===')
for shape in slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        print(f'Shape: {shape.name}')
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            for run_idx, run in enumerate(para.runs):
                if run.text.strip():
                    size = run.font.size.pt if run.font.size else 'None'
                    print(f'  [{para_idx},{run_idx}] size={size}pt text="{run.text[:50]}"')

# Processed
prs2 = Presentation('test_output.pptx')
slide2 = list(prs2.slides)[16]

print('\n=== SLIDE 17 AFTER PROCESSING ===')
for shape in slide2.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        print(f'Shape: {shape.name}')
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            for run_idx, run in enumerate(para.runs):
                if run.text.strip():
                    size = run.font.size.pt if run.font.size else 'None'
                    print(f'  [{para_idx},{run_idx}] size={size}pt text="{run.text[:50]}"')
