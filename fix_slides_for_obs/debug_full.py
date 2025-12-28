from pptx import Presentation
import fix_slides_for_obs_processor as processor

# Fresh copy
import shutil
shutil.copy('Apresentação1Original.pptx', 'test_debug.pptx')

prs = Presentation('test_debug.pptx')

# Get dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height
margin_percent = 0.05
spacing_pt = 10
margin_x = int(slide_width * margin_percent)
margin_y = int(slide_height * margin_percent)
spacing_emu = int(spacing_pt * 12700)
available_width = slide_width - (2 * margin_x)
available_height = slide_height - (2 * margin_y)
available_width_pt = available_width / 12700
available_height_pt = available_height / 12700

slide = list(prs.slides)[16]  # Slide 17

print('=== SLIDE 17 FULL DEBUG ===')
print(f'Available: {available_width_pt:.0f}pt x {available_height_pt:.0f}pt')
print()

# Collect shapes like the processor does
text_shapes = []
for shape in slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text.strip()
        
        original_sizes = {}
        min_font_in_shape = float('inf')
        max_font_in_shape = 0
        font_name = 'Arial'
        
        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            for run_idx, run in enumerate(paragraph.runs):
                if run.text.strip():
                    if run.font.size:
                        font_size = run.font.size.pt
                    else:
                        font_size = 12
                    
                    original_sizes[(para_idx, run_idx)] = font_size
                    min_font_in_shape = min(min_font_in_shape, font_size)
                    max_font_in_shape = max(max_font_in_shape, font_size)
                    
                    if run.font.name:
                        font_name = run.font.name
        
        if min_font_in_shape == float('inf'):
            min_font_in_shape = 12
        if max_font_in_shape == 0:
            max_font_in_shape = 12
        
        # Measure
        text_size = processor.measure_multiline_text_size(
            text, font_name, max_font_in_shape, available_width_pt - 20
        )
        if text_size:
            weight = text_size[1]
        else:
            weight = max_font_in_shape
        
        text_shapes.append({
            'shape': shape,
            'text': text,
            'weight': weight,
            'original_sizes': original_sizes,
            'max_font': max_font_in_shape,
            'font_name': font_name
        })
        
        print(f'Shape: {shape.name}')
        print(f'  Text: "{text[:40]}..."')
        print(f'  Max font: {max_font_in_shape}pt')
        print(f'  Weight: {weight}')
        print(f'  Original sizes: {original_sizes}')
        print()

# Calculate layout
num_shapes = len(text_shapes)
total_weight = sum(item['weight'] for item in text_shapes)
total_spacing = spacing_emu * (num_shapes - 1)
height_for_boxes = available_height - total_spacing

print(f'Total weight: {total_weight}')
print(f'Height for boxes: {height_for_boxes/12700:.0f}pt')
print()

current_y = margin_y
shape_layout = []

for item in text_shapes:
    height_ratio = item['weight'] / total_weight
    height_per_box = int(height_for_boxes * height_ratio)
    min_height = int(available_height * 0.1)
    height_per_box = max(height_per_box, min_height)
    
    shape_layout.append({
        'item': item,
        'height': height_per_box,
        'height_pt': height_per_box / 12700,
        'y': current_y
    })
    
    print(f'{item["shape"].name}:')
    print(f'  Height ratio: {height_ratio:.2%}')
    print(f'  Box height: {height_per_box/12700:.0f}pt')
    
    current_y += height_per_box + spacing_emu

print()

# Calculate scale factors
margin_pt = 10
for layout in shape_layout:
    item = layout['item']
    height_pt = layout['height_pt']
    max_font = item['max_font']
    font_name = item['font_name']
    text_content = item['text']
    
    # Binary search
    low_scale = 0.1
    high_scale = 20.0
    best_scale = 1.0
    
    for _ in range(20):
        mid_scale = (low_scale + high_scale) / 2
        max_scaled_font = max_font * mid_scale
        
        text_size = processor.measure_multiline_text_size(
            text_content, font_name, max_scaled_font, available_width_pt - margin_pt * 2
        )
        
        if text_size is None:
            high_scale = mid_scale
            continue
        
        text_width, text_height = text_size
        
        if text_height <= (height_pt - margin_pt * 2) and text_width <= (available_width_pt - margin_pt * 2):
            best_scale = mid_scale
            low_scale = mid_scale
        else:
            high_scale = mid_scale
        
        if high_scale - low_scale < 0.01:
            break
    
    layout['scale_factor'] = best_scale
    print(f'{item["shape"].name}: scale_factor = {best_scale:.2f}')
    print(f'  Original font: {max_font}pt -> Scaled: {max_font * best_scale:.0f}pt')

# Find min scale
min_scale = min(layout['scale_factor'] for layout in shape_layout)
print(f'\nMin scale factor (used for all): {min_scale:.2f}')
print()

# Apply and show results
print('=== FINAL RESULTS ===')
for layout in shape_layout:
    item = layout['item']
    for key, orig_size in item['original_sizes'].items():
        new_size = round(orig_size * min_scale)
        print(f'{item["shape"].name} {key}: {orig_size}pt -> {new_size}pt')

# Check ratio
sizes = []
for layout in shape_layout:
    item = layout['item']
    for key, orig_size in item['original_sizes'].items():
        sizes.append((item['shape'].name, orig_size, round(orig_size * min_scale)))

print()
print('=== RATIO CHECK ===')
if len(sizes) >= 2:
    print(f'Original ratio: {sizes[0][1]}:{sizes[1][1]} = {sizes[0][1]/sizes[1][1]:.3f}')
    print(f'New ratio: {sizes[0][2]}:{sizes[1][2]} = {sizes[0][2]/sizes[1][2]:.3f}')
