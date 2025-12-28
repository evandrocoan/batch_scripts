from pptx import Presentation
import fix_slides_for_obs_processor as processor

# Simulate what happens for slide 55
prs = Presentation('Apresentação1Original.pptx')
slide = list(prs.slides)[54]  # 0-indexed

slide_width = prs.slide_width
slide_height = prs.slide_height

# These values match the processor
margin_x = int(slide_width * 0.05)
margin_y = int(slide_height * 0.05)
available_width = slide_width - (2 * margin_x)
available_height = slide_height - (2 * margin_y)

available_width_pt = available_width / 12700
available_height_pt = available_height / 12700

print(f'Slide dimensions: {slide_width/12700:.0f}pt x {slide_height/12700:.0f}pt')
print(f'Available area: {available_width_pt:.0f}pt x {available_height_pt:.0f}pt')

# Simulate shape processing
for shape in slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = processor.normalize_text_whitespace(shape.text_frame.text)
        
        # Get original font info
        orig_font_size = None
        font_name = 'Arial'
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() and run.font.size:
                    orig_font_size = run.font.size.pt
                    if run.font.name:
                        font_name = run.font.name
                    break
            if orig_font_size:
                break
        
        # Calculate weight (text height at original size)
        text_size = processor.measure_multiline_text_size(
            text, font_name, orig_font_size, available_width_pt - 20
        )
        weight = text_size[1] if text_size else orig_font_size * (text.count('\n') + 1)
        
        print(f'\n=== Shape Layout Calculation ===')
        print(f'Original font: {orig_font_size}pt')
        print(f'Weight (text height at original size): {weight:.1f}pt')
        
        # For a single shape, what's the height_per_box?
        num_shapes = 1
        total_weight = weight
        spacing_emu = 0
        height_for_boxes = available_height - spacing_emu
        height_ratio = weight / total_weight  # Always 1.0 for single shape
        height_per_box = int(height_for_boxes * height_ratio)
        height_pt = height_per_box / 12700
        
        print(f'Height per box: {height_per_box} EMU = {height_pt:.1f}pt')
        
        # Now simulate binary search
        margin_pt = 10
        safety_factor = 0.98  # For single shape
        usable_height = (height_pt - margin_pt * 2) * safety_factor
        usable_width = (available_width_pt - margin_pt * 2) * safety_factor
        
        print(f'\n=== Binary Search Constraints ===')
        print(f'Usable height: {usable_height:.1f}pt')
        print(f'Usable width: {usable_width:.1f}pt')
        
        # Test different scales
        print(f'\n=== Scale Tests ===')
        for scale in [1.0, 1.5, 2.0, 2.5]:
            test_font = orig_font_size * scale
            test_size = processor.measure_multiline_text_size(
                text, font_name, test_font,
                available_width_pt - margin_pt * 2  # This is used as max_width for wrapping
            )
            if test_size:
                fits_h = test_size[1] <= usable_height
                fits_w = test_size[0] <= usable_width
                print(f'Scale {scale:.1f} ({test_font:.0f}pt): '
                      f'measured {test_size[0]:.0f}x{test_size[1]:.0f}pt, '
                      f'Height {"OK" if fits_h else f"OVERFLOW ({test_size[1]:.0f}>{usable_height:.0f})"}, '
                      f'Width {"OK" if fits_w else f"OVERFLOW ({test_size[0]:.0f}>{usable_width:.0f})"}')
