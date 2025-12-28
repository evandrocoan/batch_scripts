from pptx import Presentation
import fix_slides_for_obs_processor as processor

# Check slide 56
prs_orig = Presentation('Apresentação1Original.pptx')
prs_proc = Presentation('test_output.pptx')

slide_orig = list(prs_orig.slides)[55]  # 0-indexed
slide_proc = list(prs_proc.slides)[55]

print('=== SLIDE 56 ORIGINAL ===')
for shape in slide_orig.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        print(f'Shape: {shape.name}')
        print(f'  Position: left={shape.left/914400:.2f}" top={shape.top/914400:.2f}"')
        print(f'  Size: width={shape.width/914400:.2f}" height={shape.height/914400:.2f}"')
        text = shape.text_frame.text.strip()
        print(f'  Text: "{text[:100]}..."')
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    size = run.font.size.pt if run.font.size else 'None'
                    print(f'    Font: {size}pt - "{run.text[:50]}"')
        print()

print('\n=== SLIDE 56 AFTER PROCESSING ===')
for shape in slide_proc.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        print(f'Shape: {shape.name}')
        print(f'  Position: left={shape.left/914400:.2f}" top={shape.top/914400:.2f}"')
        print(f'  Size: width={shape.width/914400:.2f}" height={shape.height/914400:.2f}"')
        text = shape.text_frame.text.strip()
        print(f'  Text: "{text[:100]}..."')
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    size = run.font.size.pt if run.font.size else 'None'
                    print(f'    Font: {size}pt - "{run.text[:50]}"')
        print()

# Also measure the text to see if it overflows
print('\n=== TEXT MEASUREMENT ===')
slide_width = prs_proc.slide_width
slide_height = prs_proc.slide_height
available_width_pt = (slide_width - 2 * int(slide_width * 0.05)) / 12700
available_height_pt = (slide_height - 2 * int(slide_height * 0.05)) / 12700
print(f'Available area: {available_width_pt:.0f}pt x {available_height_pt:.0f}pt')

for shape in slide_proc.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text.strip()
        max_font = 0
        font_name = 'Arial'
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() and run.font.size:
                    max_font = max(max_font, run.font.size.pt)
                    if run.font.name:
                        font_name = run.font.name
        
        if max_font > 0:
            # Use the same usable_width as the processor
            margin_pt = 10
            safety_factor = 0.98
            usable_width = (available_width_pt - margin_pt * 2) * safety_factor
            usable_height = (available_height_pt - margin_pt * 2) * safety_factor
            
            text_size = processor.measure_multiline_text_size(text, font_name, max_font, usable_width)
            if text_size:
                print(f'{shape.name}: measured {text_size[0]:.0f}pt x {text_size[1]:.0f}pt at {max_font}pt font')
                box_height = shape.height / 12700
                print(f'  Box height: {box_height:.0f}pt, Text height: {text_size[1]:.0f}pt')
                print(f'  Usable height (with safety): {usable_height:.0f}pt')
                if text_size[1] > usable_height:
                    print(f'  *** OVERFLOW by {text_size[1] - usable_height:.0f}pt ***')
                else:
                    remaining = usable_height - text_size[1]
                    usage_pct = (text_size[1] / usable_height) * 100
                    print(f'  Remaining space: {remaining:.0f}pt ({usage_pct:.0f}% used)')
