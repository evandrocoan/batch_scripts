#!/usr/bin/env python3
"""
Unified slide debugging script.

This script consolidates all the individual debug scripts (check_slide*.py, debug_slide*.py, etc.)
into a single script that can debug any slide with various modes.

Usage:
    python debug_slide.py <slide_number> [options]

Examples:
    python debug_slide.py 17                    # Basic info for slide 17
    python debug_slide.py 17 --compare          # Compare original vs processed
    python debug_slide.py 17 --fonts            # Show font details
    python debug_slide.py 17 --layout           # Show layout calculation
    python debug_slide.py 17 --measurement      # Show text measurement
    python debug_slide.py 17 --shapes           # Show all shapes (including non-text)
    python debug_slide.py 17 --binary-search    # Show binary search scaling process
    python debug_slide.py 17 --lines            # Show line-by-line text analysis
    python debug_slide.py 17 --all              # Show all debug info
    python debug_slide.py --all-slides          # Show fonts for all slides
    python debug_slide.py 17 --process          # Process and show results
    python debug_slide.py --generate-output     # Process full presentation, save to tests/test_output.pptx
    python debug_slide.py --split-slides        # Split presentation into individual slide files
"""

import argparse
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Emu

# Import processor if available
try:
    import fix_slides_for_obs_processor as processor
    HAS_PROCESSOR = True
except ImportError:
    HAS_PROCESSOR = False
    print("Warning: fix_slides_for_obs_processor not found. Some features will be disabled.")

# Get the directory containing this script
SCRIPT_DIR = Path(__file__).parent.absolute()
TESTS_DIR = SCRIPT_DIR / 'tests'

# Default file paths
DEFAULT_ORIGINAL = str(SCRIPT_DIR / 'Apresentação1Original.pptx')
DEFAULT_PROCESSED = str(TESTS_DIR / 'test_output.pptx')
DEFAULT_DEBUG = str(TESTS_DIR / 'test_debug.pptx')

# Constants
EMU_PER_INCH = 914400
EMU_PER_PT = 12700


def emu_to_inches(emu):
    """Convert EMU to inches."""
    return emu / EMU_PER_INCH


def emu_to_pt(emu):
    """Convert EMU to points."""
    return emu / EMU_PER_PT


def get_slide_dimensions(prs):
    """Get slide dimensions and available area."""
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    margin_percent = 0.05
    spacing_pt = 10
    
    margin_x = int(slide_width * margin_percent)
    margin_y = int(slide_height * margin_percent)
    spacing_emu = int(spacing_pt * EMU_PER_PT)
    
    available_width = slide_width - (2 * margin_x)
    available_height = slide_height - (2 * margin_y)
    available_width_pt = emu_to_pt(available_width)
    available_height_pt = emu_to_pt(available_height)
    
    return {
        'slide_width': slide_width,
        'slide_height': slide_height,
        'margin_x': margin_x,
        'margin_y': margin_y,
        'spacing_emu': spacing_emu,
        'available_width': available_width,
        'available_height': available_height,
        'available_width_pt': available_width_pt,
        'available_height_pt': available_height_pt,
    }


def get_shape_info(shape):
    """Get information about a shape."""
    info = {
        'name': shape.name,
        'shape_type': shape.shape_type,
        'left': shape.left,
        'top': shape.top,
        'width': shape.width,
        'height': shape.height,
        'has_text_frame': shape.has_text_frame,
        'text': '',
        'runs': [],
    }
    
    if shape.has_text_frame:
        info['text'] = shape.text_frame.text.strip()
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            for run_idx, run in enumerate(para.runs):
                if run.text.strip():
                    run_info = {
                        'para_idx': para_idx,
                        'run_idx': run_idx,
                        'text': run.text,
                        'font_size': run.font.size.pt if run.font.size else None,
                        'font_name': run.font.name,
                    }
                    info['runs'].append(run_info)
    
    return info


def print_basic_info(slide, slide_num, label=""):
    """Print basic information about a slide."""
    header = f"=== SLIDE {slide_num}"
    if label:
        header += f" {label}"
    header += " ==="
    print(header)
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f'Shape: {shape.name}')
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    if run.text.strip():
                        size = run.font.size.pt if run.font.size else 'None'
                        text_preview = run.text[:50].replace('\n', '\\n')
                        print(f'  [{para_idx},{run_idx}] size={size}pt text="{text_preview}"')


def print_position_info(slide, slide_num, prs=None, label=""):
    """Print position and size information for shapes."""
    header = f"=== SLIDE {slide_num}"
    if label:
        header += f" {label}"
    header += " POSITIONS ==="
    print(header)
    
    shapes_info = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()[:60]
            shapes_info.append((shape.top, shape.name, shape, text))
    
    # Sort by top position
    shapes_info.sort(key=lambda x: x[0])
    
    for top, name, shape, text in shapes_info:
        print(f'Shape: {name}')
        print(f'  Position: left={emu_to_inches(shape.left):.2f}" top={emu_to_inches(shape.top):.2f}"')
        print(f'  Size: width={emu_to_inches(shape.width):.2f}" height={emu_to_inches(shape.height):.2f}"')
        print(f'  Text: "{text}..."')
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    size = run.font.size.pt if run.font.size else 'None'
                    print(f'    Font: {size}pt - "{run.text[:40]}"')
        print()
    
    if prs:
        print('=== SLIDE DIMENSIONS ===')
        print(f'Width: {emu_to_inches(prs.slide_width):.2f}"')
        print(f'Height: {emu_to_inches(prs.slide_height):.2f}"')


def print_all_shapes(slide, slide_num):
    """Print all shapes including non-text shapes."""
    print(f'=== SLIDE {slide_num} ALL SHAPES ===')
    
    for shape in slide.shapes:
        print(f'Shape: {shape.name}, type: {shape.shape_type}')
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print('  *** THIS IS A PICTURE ***')
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f'  Text: "{shape.text_frame.text.strip()[:50]}..."')
    
    if HAS_PROCESSOR:
        print(f'\nslide_has_visual_elements: {processor.slide_has_visual_elements(slide)}')


def print_compare(prs_orig, prs_proc, slide_idx, slide_num):
    """Compare original and processed slide."""
    slide_orig = list(prs_orig.slides)[slide_idx]
    slide_proc = list(prs_proc.slides)[slide_idx]
    
    print_position_info(slide_orig, slide_num, label="ORIGINAL")
    print()
    print_position_info(slide_proc, slide_num, label="AFTER PROCESSING")


def print_text_measurement(slide, slide_num, dims):
    """Print text measurement information."""
    if not HAS_PROCESSOR:
        print("Error: processor module required for text measurement")
        return
    
    print(f'=== SLIDE {slide_num} TEXT MEASUREMENT ===')
    print(f'Available area: {dims["available_width_pt"]:.0f}pt x {dims["available_height_pt"]:.0f}pt')
    
    margin_pt = 10
    safety_factor = 0.98
    usable_width = (dims['available_width_pt'] - margin_pt * 2) * safety_factor
    usable_height = (dims['available_height_pt'] - margin_pt * 2) * safety_factor
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()
            max_font = 0
            font_name = 'Arial'
            
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() and run.font.size:
                        max_font = max(max_font, run.font.size.pt)
                        if run.font.name:
                            font_name = run.font.name
            
            if max_font > 0:
                text_size = processor.measure_multiline_text_size(
                    text, font_name, max_font, usable_width
                )
                if text_size:
                    print(f'\n{shape.name}: measured {text_size[0]:.0f}pt x {text_size[1]:.0f}pt at {max_font}pt font')
                    box_height = emu_to_pt(shape.height)
                    print(f'  Box height: {box_height:.0f}pt, Text height: {text_size[1]:.0f}pt')
                    print(f'  Usable height (with safety): {usable_height:.0f}pt')
                    
                    if text_size[1] > usable_height:
                        print(f'  *** OVERFLOW by {text_size[1] - usable_height:.0f}pt ***')
                    else:
                        remaining = usable_height - text_size[1]
                        usage_pct = (text_size[1] / usable_height) * 100
                        print(f'  Remaining space: {remaining:.0f}pt ({usage_pct:.0f}% used)')


def print_layout_calculation(slide, slide_num, dims):
    """Print layout calculation details."""
    if not HAS_PROCESSOR:
        print("Error: processor module required for layout calculation")
        return
    
    print(f'=== SLIDE {slide_num} LAYOUT CALCULATION ===')
    print(f'Available: {dims["available_width_pt"]:.0f}pt x {dims["available_height_pt"]:.0f}pt')
    print()
    
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = processor.normalize_text_whitespace(shape.text_frame.text)
            
            original_sizes = {}
            min_font_in_shape = float('inf')
            max_font_in_shape = 0
            font_name = 'Arial'
            
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(paragraph.runs):
                    if run.text.strip():
                        font_size = run.font.size.pt if run.font.size else 12
                        original_sizes[(para_idx, run_idx)] = font_size
                        min_font_in_shape = min(min_font_in_shape, font_size)
                        max_font_in_shape = max(max_font_in_shape, font_size)
                        if run.font.name:
                            font_name = run.font.name
            
            if min_font_in_shape == float('inf'):
                min_font_in_shape = 12
            if max_font_in_shape == 0:
                max_font_in_shape = 12
            
            # Measure text
            text_size = processor.measure_multiline_text_size(
                text, font_name, max_font_in_shape, dims['available_width_pt'] - 20
            )
            weight = text_size[1] if text_size else max_font_in_shape
            
            text_shapes.append({
                'shape': shape,
                'text': text,
                'weight': weight,
                'original_sizes': original_sizes,
                'max_font': max_font_in_shape,
                'font_name': font_name
            })
            
            print(f'Shape: {shape.name}')
            print(f'  Text: "{text[:40]}..."')
            print(f'  Max font: {max_font_in_shape}pt')
            print(f'  Weight: {weight}')
            print(f'  Original sizes: {original_sizes}')
            print()
    
    if not text_shapes:
        print("No text shapes found")
        return
    
    # Calculate layout
    num_shapes = len(text_shapes)
    total_weight = sum(item['weight'] for item in text_shapes)
    total_spacing = dims['spacing_emu'] * (num_shapes - 1)
    height_for_boxes = dims['available_height'] - total_spacing
    
    print(f'Total weight: {total_weight}')
    print(f'Height for boxes: {emu_to_pt(height_for_boxes):.0f}pt')
    print()
    
    current_y = dims['margin_y']
    for item in text_shapes:
        height_ratio = item['weight'] / total_weight
        height_per_box = int(height_for_boxes * height_ratio)
        min_height = int(dims['available_height'] * 0.1)
        height_per_box = max(height_per_box, min_height)
        
        print(f'{item["shape"].name}:')
        print(f'  Height ratio: {height_ratio:.2%}')
        print(f'  Box height: {emu_to_pt(height_per_box):.0f}pt')
        current_y += height_per_box + dims['spacing_emu']


def print_binary_search(slide, slide_num, dims):
    """Print binary search scaling process."""
    if not HAS_PROCESSOR:
        print("Error: processor module required for binary search")
        return
    
    print(f'=== SLIDE {slide_num} BINARY SEARCH SCALING ===')
    print(f'Available area: {dims["available_width_pt"]:.0f}pt x {dims["available_height_pt"]:.0f}pt')
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            processor.clean_empty_paragraphs(shape.text_frame)
            text = processor.normalize_text_whitespace(shape.text_frame.text)
            
            max_font = 0
            font_name = 'Arial'
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        if run.font.size:
                            max_font = max(max_font, run.font.size.pt)
                        if run.font.name:
                            font_name = run.font.name
            
            print(f'\nShape: {shape.name}')
            print(f'Original font: {max_font}pt, font_name: {font_name}')
            print(f'Text: "{text}"')
            print(f'Text length: {len(text)} chars')
            
            margin_pt = 10
            safety_factor = 0.98
            
            height_per_box = dims['available_height']
            height_pt = emu_to_pt(height_per_box)
            
            usable_height = (height_pt - margin_pt * 2) * safety_factor
            usable_width = (dims['available_width_pt'] - margin_pt * 2) * safety_factor
            
            print(f'\nProcessor simulation:')
            print(f'  height_per_box (EMU): {height_per_box}')
            print(f'  height_pt: {height_pt:.1f}')
            print(f'  usable_height: {usable_height:.1f}')
            print(f'  usable_width: {usable_width:.1f}')
            
            # Binary search
            low_scale = 1.0
            high_scale = 30.0
            best_scale = 1.0
            
            for iteration in range(25):
                mid_scale = (low_scale + high_scale) / 2
                max_scaled_font = max_font * mid_scale
                
                text_size = processor.measure_multiline_text_size(
                    text, font_name, max_scaled_font, dims['available_width_pt'] - margin_pt * 2
                )
                
                if text_size is None:
                    all_fit = False
                else:
                    text_width, text_height = text_size
                    fits_height = text_height <= usable_height
                    fits_width = text_width <= usable_width
                    all_fit = fits_height and fits_width
                
                if iteration < 10 or not all_fit:
                    print(f'  Iter {iteration}: scale={mid_scale:.3f}, font={max_scaled_font:.1f}pt', end='')
                    if text_size:
                        print(f', size={text_size[0]:.0f}x{text_size[1]:.0f}pt', end='')
                        print(f', fits_h={fits_height}, fits_w={fits_width}', end='')
                    print(f', all_fit={all_fit}')
                
                if all_fit:
                    best_scale = mid_scale
                    low_scale = mid_scale
                else:
                    high_scale = mid_scale
                
                if high_scale - low_scale < 0.01:
                    break
            
            print(f'\nFinal best_scale: {best_scale:.3f}')
            print(f'Final font size: {max_font * best_scale:.1f}pt')


def print_line_analysis(slide, slide_num, dims):
    """Print line-by-line text analysis."""
    if not HAS_PROCESSOR:
        print("Error: processor module required for line analysis")
        return
    
    try:
        from PIL import ImageFont, ImageDraw, Image
    except ImportError:
        print("Error: PIL/Pillow required for line analysis")
        return
    
    print(f'=== SLIDE {slide_num} LINE-BY-LINE ANALYSIS ===')
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = processor.normalize_text_whitespace(shape.text_frame.text)
            
            orig_font_size = None
            font_name = 'Arial'
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() and run.font.size:
                        orig_font_size = run.font.size.pt
                        if run.font.name:
                            font_name = run.font.name
                        break
                if orig_font_size:
                    break
            
            if not orig_font_size:
                orig_font_size = 12
            
            font_path = processor.get_font_path(font_name)
            
            print(f'\nShape: {shape.name}')
            print(f'Text (total {len(text)} chars):')
            print(f'"{text}"')
            print()
            
            lines = text.split('\n')
            print(f'Number of lines: {len(lines)}')
            print()
            
            margin_pt = 10
            usable_width = dims['available_width_pt'] - margin_pt * 2
            
            for scale in [1.0, 1.5, 2.0]:
                test_font_pt = orig_font_size * scale
                try:
                    font = ImageFont.truetype(font_path, int(test_font_pt))
                    img = Image.new('RGB', (1, 1))
                    draw = ImageDraw.Draw(img)
                    
                    print(f'=== Scale {scale:.1f} ({test_font_pt:.0f}pt) ===')
                    for i, line in enumerate(lines):
                        bbox = draw.textbbox((0, 0), line, font=font)
                        width = bbox[2] - bbox[0]
                        line_preview = line[:60] if len(line) > 60 else line
                        print(f'  Line {i+1} ({len(line)} chars): {width:.0f}pt wide - "{line_preview}..."')
                    
                    text_size = processor.measure_multiline_text_size(
                        text, font_name, test_font_pt, usable_width
                    )
                    if text_size:
                        print(f'  -> After wrapping to {usable_width:.0f}pt: {text_size[0]:.0f}x{text_size[1]:.0f}pt')
                    print()
                except Exception as e:
                    print(f'Error loading font: {e}')


def print_all_slides_fonts(prs):
    """Print font information for all slides."""
    print('=== ALL SLIDES FONT SIZES ===')
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f'\nSlide {slide_num}:')
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                print(f'  Shape: {shape.name}')
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(para.runs):
                        if run.text.strip():
                            size = run.font.size.pt if run.font.size else 'None'
                            text_preview = run.text[:30].replace('\n', '\\n')
                            print(f'    [{para_idx},{run_idx}] "{text_preview}" -> {size}pt')


def process_and_show_results(prs_path, slide_num):
    """Process a presentation and show results."""
    if not HAS_PROCESSOR:
        print("Error: processor module required for processing")
        return
    
    # Make a copy
    debug_path = DEFAULT_DEBUG
    shutil.copy(prs_path, debug_path)
    
    prs = Presentation(debug_path)
    
    print('=== PROCESSING WITH REPOSITION ===')
    result = processor.reposition_and_maximize_font(prs, margin_percent=0.05, spacing_pt=10)
    
    print(f'\nSlides processed: {result["slides_processed"]}')
    print(f'\nFont changes for slide {slide_num}:')
    
    for change in result['font_changes']:
        if change['slide_num'] == slide_num:
            scale = change.get('scale', 'N/A')
            scale_str = f'{scale:.2f}' if isinstance(scale, (int, float)) else str(scale)
            print(f"  {change['old_size']}pt -> {change['new_size']}pt (scale: {scale_str})")
    
    # Save and show final state
    prs.save(debug_path)
    
    print(f'\n=== SLIDE {slide_num} FINAL STATE ===')
    prs2 = Presentation(debug_path)
    slide_idx = slide_num - 1
    
    if slide_idx < len(list(prs2.slides)):
        slide = list(prs2.slides)[slide_idx]
        print_position_info(slide, slide_num, prs2)


def print_scale_tests(slide, slide_num, dims):
    """Print scale factor tests."""
    if not HAS_PROCESSOR:
        print("Error: processor module required for scale tests")
        return
    
    print(f'=== SLIDE {slide_num} SCALE TESTS ===')
    print(f'Available area: {dims["available_width_pt"]:.0f}pt x {dims["available_height_pt"]:.0f}pt')
    
    margin_pt = 10
    safety_factor = 0.98
    usable_height = (dims['available_height_pt'] - margin_pt * 2) * safety_factor
    usable_width = (dims['available_width_pt'] - margin_pt * 2) * safety_factor
    
    print(f'Usable area: {usable_width:.0f}pt x {usable_height:.0f}pt')
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = processor.normalize_text_whitespace(shape.text_frame.text)
            
            max_font = 0
            font_name = 'Arial'
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        if run.font.size:
                            max_font = max(max_font, run.font.size.pt)
                        if run.font.name:
                            font_name = run.font.name
            
            if max_font == 0:
                max_font = 12
            
            print(f'\nShape: {shape.name}')
            print(f'Original font: {max_font}pt, font_name: {font_name}')
            print(f'Text: "{text}"')
            print('\nScale tests:')
            
            for scale in [1.0, 1.5, 2.0, 2.5, 3.0, 3.5, 4.0]:
                scaled_font = max_font * scale
                text_size = processor.measure_multiline_text_size(
                    text, font_name, scaled_font, dims['available_width_pt'] - margin_pt * 2
                )
                if text_size:
                    fits_h = "OK" if text_size[1] <= usable_height else "OVERFLOW"
                    fits_w = "OK" if text_size[0] <= usable_width else "OVERFLOW"
                    print(f'  Scale {scale:.1f}: font={scaled_font:.0f}pt, '
                          f'size={text_size[0]:.0f}x{text_size[1]:.0f}pt, '
                          f'height:{fits_h}, width:{fits_w}')


def generate_test_output(input_path):
    """Process a presentation and save to test_output.pptx."""
    if not HAS_PROCESSOR:
        print("Error: processor module required for generating output")
        return
    
    if not Path(input_path).exists():
        print(f"Error: File not found: {input_path}")
        return
    
    prs = Presentation(input_path)
    result = processor.reposition_and_maximize_font(prs)
    print(f'Slides processed: {result["slides_processed"]}')
    
    output_path = DEFAULT_PROCESSED
    prs.save(output_path)
    print(f'Saved to {output_path}')


def delete_slide(prs, slide_index):
    """Delete a slide from presentation by index."""
    slide_id = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[slide_index]


def extract_single_slide(source_path, slide_index, output_path):
    """
    Extract a single slide from presentation while preserving exact structure.
    
    This works by loading the full presentation and deleting all other slides.
    This preserves the original shape types (Title, Subtitle, etc.)
    """
    prs = Presentation(source_path)
    total_slides = len(prs.slides)
    
    # Delete slides from end to beginning to maintain indices
    # Keep only the slide at slide_index
    for i in range(total_slides - 1, -1, -1):
        if i != slide_index:
            delete_slide(prs, i)
    
    prs.save(output_path)


def split_presentation(source_path):
    """Split presentation into individual slide files."""
    import os
    
    if not Path(source_path).exists():
        print(f"Error: File not found: {source_path}")
        return
    
    # Create output directory in tests/
    output_dir = TESTS_DIR / 'test_slides'
    os.makedirs(output_dir, exist_ok=True)
    
    # Get total slide count
    prs = Presentation(source_path)
    total_slides = len(prs.slides)
    print(f"Splitting {total_slides} slides from {source_path}...")
    
    for i in range(total_slides):
        slide_num = i + 1
        output_path = os.path.join(output_dir, f'slide_{slide_num:02d}.pptx')
        
        # Extract single slide (preserves original shape types)
        extract_single_slide(source_path, i, output_path)
        print(f"  Saved: {output_path}")
    
    print(f"\nDone! Created {total_slides} individual slide files in '{output_dir}/'")


def main():
    parser = argparse.ArgumentParser(
        description='Unified slide debugging script',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    
    parser.add_argument('slide', type=int, nargs='?', default=None,
                        help='Slide number (1-indexed)')
    parser.add_argument('--original', '-o', default=DEFAULT_ORIGINAL,
                        help=f'Original presentation file (default: {DEFAULT_ORIGINAL})')
    parser.add_argument('--processed', '-p', default=DEFAULT_PROCESSED,
                        help=f'Processed presentation file (default: {DEFAULT_PROCESSED})')
    
    # Debug modes
    parser.add_argument('--compare', '-c', action='store_true',
                        help='Compare original vs processed')
    parser.add_argument('--fonts', '-f', action='store_true',
                        help='Show font details')
    parser.add_argument('--layout', '-l', action='store_true',
                        help='Show layout calculation')
    parser.add_argument('--measurement', '-m', action='store_true',
                        help='Show text measurement')
    parser.add_argument('--shapes', '-s', action='store_true',
                        help='Show all shapes (including non-text)')
    parser.add_argument('--binary-search', '-b', action='store_true',
                        help='Show binary search scaling process')
    parser.add_argument('--lines', action='store_true',
                        help='Show line-by-line text analysis')
    parser.add_argument('--scale-tests', '-t', action='store_true',
                        help='Show scale factor tests')
    parser.add_argument('--process', action='store_true',
                        help='Process and show results')
    parser.add_argument('--all', '-a', action='store_true',
                        help='Show all debug info')
    parser.add_argument('--all-slides', action='store_true',
                        help='Show fonts for all slides')
    parser.add_argument('--generate-output', action='store_true',
                        help='Process full presentation and save to test_output.pptx')
    parser.add_argument('--split-slides', action='store_true',
                        help='Split presentation into individual slide files in tests/test_slides/')
    
    args = parser.parse_args()
    
    # Check for generate-output mode
    if args.generate_output:
        generate_test_output(args.original)
        return
    
    # Check for split-slides mode
    if args.split_slides:
        split_presentation(args.original)
        return
    
    # Check for all-slides mode
    if args.all_slides:
        if not Path(args.original).exists():
            print(f"Error: File not found: {args.original}")
            sys.exit(1)
        prs = Presentation(args.original)
        print_all_slides_fonts(prs)
        return
    
    # Require slide number for other modes
    if args.slide is None:
        parser.print_help()
        print("\nError: Slide number is required (unless using --all-slides, --generate-output, or --split-slides)")
        sys.exit(1)
    
    slide_num = args.slide
    slide_idx = slide_num - 1
    
    # Check files exist
    if not Path(args.original).exists():
        print(f"Error: Original file not found: {args.original}")
        sys.exit(1)
    
    # Load presentations
    prs_orig = Presentation(args.original)
    
    # Check slide number is valid
    total_slides = len(list(prs_orig.slides))
    if slide_idx < 0 or slide_idx >= total_slides:
        print(f"Error: Slide {slide_num} out of range. Presentation has {total_slides} slides.")
        sys.exit(1)
    
    slide_orig = list(prs_orig.slides)[slide_idx]
    dims = get_slide_dimensions(prs_orig)
    
    # Determine what to show
    show_basic = not any([args.compare, args.fonts, args.layout, args.measurement,
                          args.shapes, args.binary_search, args.lines, args.scale_tests,
                          args.process, args.all])
    
    if args.all or show_basic:
        print_basic_info(slide_orig, slide_num, "ORIGINAL")
        print()
        print_position_info(slide_orig, slide_num, prs_orig, "ORIGINAL")
    
    if args.compare or args.all:
        if Path(args.processed).exists():
            prs_proc = Presentation(args.processed)
            print()
            print_compare(prs_orig, prs_proc, slide_idx, slide_num)
        else:
            print(f"\nWarning: Processed file not found: {args.processed}")
    
    if args.fonts or args.all:
        print()
        print_basic_info(slide_orig, slide_num, "FONTS")
    
    if args.shapes or args.all:
        print()
        print_all_shapes(slide_orig, slide_num)
    
    if args.layout or args.all:
        print()
        print_layout_calculation(slide_orig, slide_num, dims)
    
    if args.measurement or args.all:
        print()
        print_text_measurement(slide_orig, slide_num, dims)
    
    if args.scale_tests or args.all:
        print()
        print_scale_tests(slide_orig, slide_num, dims)
    
    if args.binary_search or args.all:
        print()
        print_binary_search(slide_orig, slide_num, dims)
    
    if args.lines or args.all:
        print()
        print_line_analysis(slide_orig, slide_num, dims)
    
    if args.process:
        print()
        process_and_show_results(args.original, slide_num)


if __name__ == '__main__':
    main()
