"""
Split the original presentation into 72 individual slide files.
Each file contains just one slide for isolated testing.

This method preserves exact slide structure by:
1. Loading the original presentation for each slide
2. Deleting all slides except the target one
3. Saving as new file

This ensures shapes maintain their original types (Title, Subtitle, etc.)
"""
from pptx import Presentation
import os


def delete_slide(prs, slide_index):
    """Delete a slide from presentation by index."""
    slide_id = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[slide_index]


def extract_single_slide(source_path, slide_index, output_path):
    """
    Extract a single slide from presentation while preserving exact structure.
    
    This works by loading the full presentation and deleting all other slides.
    This preserves the original shape types (Title, Subtitle, etc.)
    """
    # Load fresh copy of presentation
    prs = Presentation(source_path)
    
    total_slides = len(prs.slides)
    
    # Delete slides from end to beginning to maintain indices
    # Keep only the slide at slide_index
    for i in range(total_slides - 1, -1, -1):
        if i != slide_index:
            delete_slide(prs, i)
    
    # Save single slide presentation
    prs.save(output_path)


def split_presentation():
    """Split presentation into individual slide files."""
    # Create output directory
    output_dir = 'test_slides'
    os.makedirs(output_dir, exist_ok=True)
    
    # Source presentation
    source_path = 'Apresentação1Original.pptx'
    
    # Get total slide count
    prs = Presentation(source_path)
    total_slides = len(prs.slides)
    print(f"Splitting {total_slides} slides from {source_path}...")
    
    for i in range(total_slides):
        slide_num = i + 1
        output_path = os.path.join(output_dir, f'slide_{slide_num:02d}.pptx')
        
        # Extract single slide (preserves original shape types)
        extract_single_slide(source_path, i, output_path)
        print(f"  Saved: {output_path}")
    
    print(f"\nDone! Created {total_slides} individual slide files in '{output_dir}/'")


if __name__ == '__main__':
    split_presentation()
