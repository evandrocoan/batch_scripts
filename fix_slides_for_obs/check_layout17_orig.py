from pptx import Presentation
from pptx.util import Pt, Emu

# Check ORIGINAL positions and sizes of slide 17
prs = Presentation('Apresentação1Original.pptx')
slide = list(prs.slides)[16]

print('=== SLIDE 17 ORIGINAL SHAPE POSITIONS ===')
for shape in slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        print(f'Shape: {shape.name}')
        print(f'  Position: left={shape.left/914400:.2f}" top={shape.top/914400:.2f}"')
        print(f'  Size: width={shape.width/914400:.2f}" height={shape.height/914400:.2f}"')
        print(f'  Text: "{shape.text_frame.text[:50]}..."')
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    size = run.font.size.pt if run.font.size else 'None'
                    print(f'  Font: {size}pt')
        print()

print('=== SLIDE DIMENSIONS ===')
print(f'Width: {prs.slide_width/914400:.2f}"')
print(f'Height: {prs.slide_height/914400:.2f}"')
