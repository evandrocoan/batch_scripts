from pptx import Presentation
import fix_slides_for_obs_processor as processor

# Fresh copy
import shutil
shutil.copy('Apresentação1Original.pptx', 'test_debug.pptx')

prs = Presentation('test_debug.pptx')
slide = list(prs.slides)[16]  # Slide 17

print('=== SLIDE 17 DEBUG ===')
print()

# Get slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height
margin_percent = 0.05
margin_x = int(slide_width * margin_percent)
margin_y = int(slide_height * margin_percent)
available_width = slide_width - (2 * margin_x)
available_height = slide_height - (2 * margin_y)
available_width_pt = available_width / 12700
available_height_pt = available_height / 12700

print(f'Available width: {available_width_pt:.0f}pt')
print(f'Available height: {available_height_pt:.0f}pt')
print()

for shape in slide.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text.strip()
        
        # Get max font
        max_font = 0
        font_name = 'Arial'
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() and run.font.size:
                    max_font = max(max_font, run.font.size.pt)
                    if run.font.name:
                        font_name = run.font.name
        
        if max_font == 0:
            max_font = 12
        
        # Measure text
        text_size = processor.measure_multiline_text_size(
            text, font_name, max_font, available_width_pt - 20
        )
        
        print(f'Shape: {shape.name}')
        print(f'  Text: "{text[:50]}..."')
        print(f'  Original max font: {max_font}pt')
        print(f'  Font name: {font_name}')
        if text_size:
            print(f'  Measured size: {text_size[0]:.0f}pt x {text_size[1]:.0f}pt (width x height)')
            print(f'  Weight (height): {text_size[1]:.0f}')
        else:
            print(f'  Measured size: FAILED')
        print()
