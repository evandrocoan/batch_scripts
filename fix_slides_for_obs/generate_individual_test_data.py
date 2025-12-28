"""
Generate expected test data by processing each individual slide file.
This captures the actual results of processing each slide in isolation.
"""
import os
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fix_slides_for_obs_processor as processor

def get_text_shapes(slide):
    """Get all text shapes from a slide."""
    shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            if shape.text_frame.text.strip() != '[IMAGE]':
                shapes.append(shape)
    return shapes

def get_max_font_size(slide):
    """Get maximum font size in slide."""
    max_font = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_font = max(max_font, run.font.size.pt)
    return max_font if max_font > 0 else None

def slide_has_visual_elements(slide):
    """Check if slide has visual elements."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == '[IMAGE]':
            return True
    return False

def generate_expected_results():
    """Process each slide individually and capture results."""
    results = {}
    
    for slide_num in range(1, 73):
        input_file = f'test_slides/slide_{slide_num:02d}.pptx'
        
        if not os.path.exists(input_file):
            print(f"Slide {slide_num:02d}: File not found")
            continue
        
        # Load original
        prs_orig = Presentation(input_file)
        slide_orig = prs_orig.slides[0]
        
        has_visual = slide_has_visual_elements(slide_orig)
        orig_shapes = get_text_shapes(slide_orig)
        orig_font = get_max_font_size(slide_orig)
        has_text = len(orig_shapes) > 0
        
        # Process
        prs_proc = Presentation(input_file)
        processor.reposition_and_maximize_font(prs_proc)
        
        # Save to temp file and reload
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            temp_path = f.name
        prs_proc.save(temp_path)
        
        prs_result = Presentation(temp_path)
        slide_result = prs_result.slides[0]
        
        result_shapes = get_text_shapes(slide_result)
        result_font = get_max_font_size(slide_result)
        
        # Determine if transformed
        was_transformed = False
        if orig_shapes and result_shapes:
            if orig_shapes[0].left != result_shapes[0].left or \
               orig_shapes[0].top != result_shapes[0].top:
                was_transformed = True
            if orig_font and result_font and abs(orig_font - result_font) > 0.5:
                was_transformed = True
        
        should_transform = not has_visual and has_text
        
        result = {
            'should_transform': should_transform and was_transformed,
            'has_visual': has_visual,
        }
        
        if not has_text:
            result['has_text'] = False
        
        if orig_font:
            result['orig_font'] = round(orig_font, 1)
        if result_font:
            result['proc_font'] = round(result_font, 1)
        
        results[slide_num] = result
        
        # Print progress
        status = "TRANSFORM" if was_transformed else ("VISUAL" if has_visual else "UNCHANGED")
        font_info = f" ({orig_font:.1f}pt -> {result_font:.1f}pt)" if orig_font and result_font else ""
        print(f"Slide {slide_num:02d}: {status}{font_info}")
        
        # Cleanup
        os.unlink(temp_path)
    
    return results

def print_python_dict(results):
    """Print results as Python dictionary for copy-paste."""
    print("\n\n# Expected results dictionary:")
    print("EXPECTED_RESULTS = {")
    for slide_num in sorted(results.keys()):
        r = results[slide_num]
        parts = []
        parts.append(f"'should_transform': {r['should_transform']}")
        parts.append(f"'has_visual': {r['has_visual']}")
        if 'has_text' in r:
            parts.append(f"'has_text': {r['has_text']}")
        if 'orig_font' in r:
            parts.append(f"'orig_font': {r['orig_font']}")
        if 'proc_font' in r:
            parts.append(f"'proc_font': {r['proc_font']}")
        
        print(f"    {slide_num}: {{{', '.join(parts)}}},")
    print("}")


if __name__ == '__main__':
    results = generate_expected_results()
    print_python_dict(results)
