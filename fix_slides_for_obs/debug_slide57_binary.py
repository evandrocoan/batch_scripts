from pptx import Presentation
import fix_slides_for_obs_processor as processor

# Debug slide 57 scaling - more detailed
prs_orig = Presentation('Apresentação1Original.pptx')

slide_orig = list(prs_orig.slides)[56]  # 0-indexed

slide_width = prs_orig.slide_width
slide_height = prs_orig.slide_height

margin_percent = 0.05
margin_x = int(slide_width * margin_percent)
margin_y = int(slide_height * margin_percent)

available_width = slide_width - (2 * margin_x)
available_height = slide_height - (2 * margin_y)
available_width_pt = available_width / 12700
available_height_pt = available_height / 12700

print(f'Slide dimensions: {slide_width/914400:.2f}" x {slide_height/914400:.2f}"')
print(f'Available area: {available_width_pt:.0f}pt x {available_height_pt:.0f}pt')

for shape in slide_orig.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        # Clean empty paragraphs first (like the processor does)
        processor.clean_empty_paragraphs(shape.text_frame)
        text = processor.normalize_text_whitespace(shape.text_frame.text)
        
        # Get original font info
        max_font = 0
        font_name = 'Arial'
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    if run.font.size:
                        max_font = max(max_font, run.font.size.pt)
                    if run.font.name:
                        font_name = run.font.name
        
        print(f'\nShape: {shape.name}')
        print(f'Original font: {max_font}pt, font_name: {font_name}')
        print(f'Text after normalization: "{text}"')
        print(f'Text length: {len(text)} chars')
        
        # Simulate what the processor does - use SAME values as processor
        margin_pt = 10  # Same as processor
        safety_factor = 0.98  # Same as processor for single shape
        
        # This is height_per_box = available_height (since single shape, ratio=1.0)
        height_per_box = available_height
        height_pt = height_per_box / 12700
        
        usable_height = (height_pt - margin_pt * 2) * safety_factor
        usable_width = (available_width_pt - margin_pt * 2) * safety_factor
        
        print(f'\nProcessor simulation:')
        print(f'  height_per_box (EMU): {height_per_box}')
        print(f'  height_pt: {height_pt:.1f}')
        print(f'  usable_height: {usable_height:.1f}')
        print(f'  usable_width: {usable_width:.1f}')
        
        # Binary search like processor
        low_scale = 1.0
        high_scale = 30.0
        best_scale = 1.0
        
        for iteration in range(25):
            mid_scale = (low_scale + high_scale) / 2
            max_scaled_font = max_font * mid_scale
            
            text_size = processor.measure_multiline_text_size(
                text, font_name, max_scaled_font, available_width_pt - margin_pt * 2
            )
            
            if text_size is None:
                all_fit = False
            else:
                text_width, text_height = text_size
                fits_height = text_height <= usable_height
                fits_width = text_width <= usable_width
                all_fit = fits_height and fits_width
            
            if iteration < 10 or not all_fit:
                print(f'  Iter {iteration}: scale={mid_scale:.3f}, font={max_scaled_font:.1f}pt', end='')
                if text_size:
                    print(f', size={text_size[0]:.0f}x{text_size[1]:.0f}pt', end='')
                    print(f', fits_h={fits_height}, fits_w={fits_width}', end='')
                print(f', all_fit={all_fit}')
            
            if all_fit:
                best_scale = mid_scale
                low_scale = mid_scale
            else:
                high_scale = mid_scale
            
            if high_scale - low_scale < 0.01:
                break
        
        print(f'\nFinal best_scale: {best_scale:.3f}')
        print(f'Final font size: {max_font * best_scale:.1f}pt')
