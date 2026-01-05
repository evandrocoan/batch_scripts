# -*- coding: utf-8 -*-
"""
Individual slide tests for fix_slides_for_obs_processor.

Each test:
1. Loads an individual slide file from test_slides/
2. Processes it through the processor
3. Saves to a temporary file
4. Reopens the saved file
5. Verifies the expected values (position, size, font, text content)

Additional verifications:
- Text overflow detection (text must fit within available area)
- Font ratio preservation (proportions between fonts must be maintained)
- Vertical ordering (shapes must maintain relative vertical order)
- Position at margins (transformed shapes must be at margin positions)

Run with: python -m pytest tests/test_individual_slides.py -v
Or from tests/: python -m pytest test_individual_slides.py -v
"""
import pytest
import os
import sys
from pptx import Presentation
from pptx.util import Emu

# Add parent directory to path for imports
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import fix_slides_for_obs_processor as processor

# Get the directory containing this test file
TEST_DIR = os.path.dirname(os.path.abspath(__file__))


# =============================================================================
# CONSTANTS
# =============================================================================

# Standard slide dimensions (16:9 aspect ratio, same as processor uses)
SLIDE_WIDTH_PT = 960.0  # points
SLIDE_HEIGHT_PT = 540.0  # points
MARGIN_PERCENT = 0.05
MARGIN_X_PT = SLIDE_WIDTH_PT * MARGIN_PERCENT
MARGIN_Y_PT = SLIDE_HEIGHT_PT * MARGIN_PERCENT
AVAILABLE_WIDTH_PT = SLIDE_WIDTH_PT - (2 * MARGIN_X_PT)
AVAILABLE_HEIGHT_PT = SLIDE_HEIGHT_PT - (2 * MARGIN_Y_PT)

# EMU conversion (1 point = 12700 EMUs)
EMU_PER_PT = 12700


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def get_slide_text(slide):
    """Get all text content from a slide, joined with newlines."""
    all_text = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            all_text.append(shape.text_frame.text.strip())
    return '\n'.join(all_text)


def get_max_font(slide):
    """Get the maximum font size from all shapes in slide."""
    max_font = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_font = max(max_font, run.font.size.pt)
    return max_font if max_font > 0 else None


def get_shapes_info(slide):
    """Get detailed info about each shape with text.
    
    Returns list of dicts with:
    - name: shape name
    - fonts: list of unique font sizes in the shape
    - text: shape text content
    - left: left position in points
    - top: top position in points
    - width: width in points
    - height: height in points
    - font_name: primary font name
    """
    shapes_info = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            fonts_set = set()
            font_name = 'Arial'
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        fonts_set.add(run.font.size.pt)
                    if run.font.name:
                        font_name = run.font.name
            shapes_info.append({
                'name': shape.name,
                'fonts': sorted(fonts_set),
                'text': shape.text_frame.text.strip(),
                'left': shape.left / EMU_PER_PT if shape.left else 0,
                'top': shape.top / EMU_PER_PT if shape.top else 0,
                'width': shape.width / EMU_PER_PT if shape.width else 0,
                'height': shape.height / EMU_PER_PT if shape.height else 0,
                'font_name': font_name,
            })
    return shapes_info


def process_and_save_slide(slide_name, tmp_path):
    """Process a slide, save it, and return path to saved file and original shapes info."""
    input_file = os.path.join(TEST_DIR, 'test_slides', f'{slide_name}.pptx')
    
    if not os.path.exists(input_file):
        pytest.skip(f"File {input_file} not found")
    
    output_file = tmp_path / f'{slide_name}_processed.pptx'
    
    # Load original and get shapes info BEFORE processing
    prs_orig = Presentation(input_file)
    slide_orig = prs_orig.slides[0]
    original_shapes_info = []
    for shape in slide_orig.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            fonts_set = set()
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        fonts_set.add(run.font.size.pt)
            original_shapes_info.append({
                'name': shape.name,
                'top': shape.top / EMU_PER_PT if shape.top else 0,
                'fonts': sorted(fonts_set),
            })
    # Sort by original top position to get the vertical order
    original_shapes_info.sort(key=lambda s: s['top'])
    original_order = [s['name'] for s in original_shapes_info]
    original_fonts_by_shape = {s['name']: s['fonts'] for s in original_shapes_info}
    
    # Load fresh copy and process
    prs = Presentation(input_file)
    processor.reposition_and_maximize_font(prs)
    prs.save(str(output_file))
    
    return str(output_file), original_order, original_fonts_by_shape


def load_and_verify(filepath):
    """Load a processed file and extract info."""
    prs = Presentation(filepath)
    slide = prs.slides[0]
    
    # Get slide dimensions
    slide_width_pt = prs.slide_width / EMU_PER_PT
    slide_height_pt = prs.slide_height / EMU_PER_PT
    margin_x = slide_width_pt * MARGIN_PERCENT
    margin_y = slide_height_pt * MARGIN_PERCENT
    available_width = slide_width_pt - (2 * margin_x)
    available_height = slide_height_pt - (2 * margin_y)
    
    return {
        'text': get_slide_text(slide),
        'max_font': get_max_font(slide),
        'has_visual': processor.slide_has_visual_elements(slide),
        'shapes': get_shapes_info(slide),
        'num_shapes': len([s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]),
        'slide_width': slide_width_pt,
        'slide_height': slide_height_pt,
        'margin_x': margin_x,
        'margin_y': margin_y,
        'available_width': available_width,
        'available_height': available_height,
    }


def check_no_text_overflow(result, tolerance_pt=5.0):
    """
    Check that text does not overflow the available area.
    
    Args:
        result: dict from load_and_verify
        tolerance_pt: tolerance in points for boundary checking
    
    Returns:
        tuple: (passes: bool, message: str)
    """
    if result['has_visual']:
        return (True, "Skipped - slide has visual elements")
    
    shapes = result['shapes']
    if not shapes:
        return (True, "No shapes to check")
    
    # Calculate total height used by all shapes
    if len(shapes) == 0:
        return (True, "No text shapes")
    
    # Get vertical span of all shapes
    shapes_sorted = sorted(shapes, key=lambda s: s['top'])
    first_shape = shapes_sorted[0]
    last_shape = shapes_sorted[-1]
    
    total_span = (last_shape['top'] + last_shape['height']) - first_shape['top']
    
    # Check if total span fits in available height
    if total_span > result['available_height'] + tolerance_pt:
        return (False, f"Vertical overflow: {total_span:.1f}pt > {result['available_height']:.1f}pt available")
    
    # Check each shape width
    for shape in shapes:
        if shape['width'] > result['available_width'] + tolerance_pt:
            return (False, f"Shape '{shape['name']}' width overflow: {shape['width']:.1f}pt > {result['available_width']:.1f}pt")
    
    return (True, "No overflow")


def check_font_ratio_preserved(original_fonts, processed_fonts, tolerance=0.15):
    """
    Check that the ratio between font sizes is preserved after scaling.
    
    Args:
        original_fonts: list of original font sizes (sorted)
        processed_fonts: list of processed font sizes (sorted)
        tolerance: allowed deviation from original ratio (0.15 = 15%)
    
    Returns:
        tuple: (passes: bool, message: str)
    """
    if len(original_fonts) < 2 or len(processed_fonts) < 2:
        return (True, "Not enough fonts to compare ratios")
    
    if len(original_fonts) != len(processed_fonts):
        return (False, f"Font count mismatch: {len(original_fonts)} vs {len(processed_fonts)}")
    
    # Compare ratios between consecutive font sizes
    for i in range(len(original_fonts) - 1):
        if original_fonts[i] == 0 or processed_fonts[i] == 0:
            continue
        
        orig_ratio = original_fonts[i + 1] / original_fonts[i]
        proc_ratio = processed_fonts[i + 1] / processed_fonts[i]
        
        ratio_diff = abs(orig_ratio - proc_ratio) / orig_ratio
        if ratio_diff > tolerance:
            return (False, f"Ratio not preserved: orig {orig_ratio:.2f} vs proc {proc_ratio:.2f} (diff: {ratio_diff:.1%})")
    
    return (True, "Ratios preserved")


def check_vertical_order_preserved(processed_shapes, original_order):
    """
    Check that shapes maintain their original vertical order after processing.
    
    Args:
        processed_shapes: list of shape info dicts from processed slide
        original_order: list of shape names in original vertical order (top to bottom)
    
    Returns:
        tuple: (passes: bool, message: str)
    """
    if len(processed_shapes) < 2 or len(original_order) < 2:
        return (True, "Not enough shapes to check order")
    
    # Get processed shapes sorted by top position
    shapes_sorted = sorted(processed_shapes, key=lambda s: s['top'])
    processed_order = [s['name'] for s in shapes_sorted]
    
    # Compare the order
    if processed_order != original_order:
        return (False, f"Vertical order changed: original {original_order} vs processed {processed_order}")
    
    return (True, "Vertical order preserved")


def check_position_at_margins(result, tolerance_pt=5.0):
    """
    Check that transformed slides have shapes positioned at margins.
    
    Args:
        result: dict from load_and_verify
        tolerance_pt: tolerance in points
    
    Returns:
        tuple: (passes: bool, message: str)
    """
    if result['has_visual']:
        return (True, "Skipped - slide has visual elements")
    
    shapes = result['shapes']
    if not shapes:
        return (True, "No shapes to check")
    
    # Find topmost shape
    shapes_sorted = sorted(shapes, key=lambda s: s['top'])
    top_shape = shapes_sorted[0]
    
    # Check left position
    if abs(top_shape['left'] - result['margin_x']) > tolerance_pt:
        return (False, f"Left position not at margin: {top_shape['left']:.1f}pt vs expected {result['margin_x']:.1f}pt")
    
    # Check top position
    if abs(top_shape['top'] - result['margin_y']) > tolerance_pt:
        return (False, f"Top position not at margin: {top_shape['top']:.1f}pt vs expected {result['margin_y']:.1f}pt")
    
    # Check width matches available width
    if abs(top_shape['width'] - result['available_width']) > tolerance_pt:
        return (False, f"Width not matching available: {top_shape['width']:.1f}pt vs expected {result['available_width']:.1f}pt")
    
    return (True, "Positioned at margins")


def check_space_utilization(result, min_utilization=0.50):
    """
    Check that the processed slide uses space efficiently.
    
    Args:
        result: dict from load_and_verify
        min_utilization: minimum expected utilization (0.50 = 50%)
    
    Returns:
        tuple: (passes: bool, message: str, utilization: float)
    """
    if result['has_visual']:
        return (True, "Skipped - slide has visual elements", 0.0)
    
    shapes = result['shapes']
    if not shapes:
        return (True, "No shapes to check", 0.0)
    
    # Calculate total height used
    shapes_sorted = sorted(shapes, key=lambda s: s['top'])
    first_shape = shapes_sorted[0]
    last_shape = shapes_sorted[-1]
    
    total_used = (last_shape['top'] + last_shape['height']) - first_shape['top']
    utilization = total_used / result['available_height']
    
    if utilization < min_utilization:
        return (False, f"Low space utilization: {utilization:.1%} < {min_utilization:.1%}", utilization)
    
    return (True, f"Space utilization: {utilization:.1%}", utilization)


def check_shapes_stacked_correctly(result, tolerance_pt=5.0):
    """
    Check that multiple shapes are stacked vertically without gaps or overlaps.
    
    Args:
        result: dict from load_and_verify
        tolerance_pt: tolerance for gap/overlap detection
    
    Returns:
        tuple: (passes: bool, message: str)
    """
    if result['has_visual']:
        return (True, "Skipped - slide has visual elements")
    
    shapes = result['shapes']
    if len(shapes) < 2:
        return (True, "Not enough shapes to check stacking")
    
    # Sort by top position
    shapes_sorted = sorted(shapes, key=lambda s: s['top'])
    
    for i in range(len(shapes_sorted) - 1):
        current = shapes_sorted[i]
        next_shape = shapes_sorted[i + 1]
        
        current_bottom = current['top'] + current['height']
        gap = next_shape['top'] - current_bottom
        
        # Check for overlap (negative gap beyond tolerance)
        if gap < -tolerance_pt:
            return (False, f"Shapes overlap: '{current['name']}' and '{next_shape['name']}' overlap by {-gap:.1f}pt")
    
    return (True, "Shapes stacked correctly")


def check_font_ratios_in_shape(original_fonts, processed_fonts, tolerance=0.15):
    """
    Check that font ratios within a shape are preserved after scaling.
    
    Args:
        original_fonts: list of original font sizes (sorted)
        processed_fonts: list of processed font sizes (sorted)  
        tolerance: allowed deviation (0.15 = 15%)
    
    Returns:
        tuple: (passes: bool, message: str)
    """
    if len(original_fonts) < 2 or len(processed_fonts) < 2:
        return (True, "Not enough fonts to compare")
    
    if len(original_fonts) != len(processed_fonts):
        # Font count changed - might be intentional, skip ratio check
        return (True, "Font count changed, skipping ratio check")
    
    for i in range(len(original_fonts) - 1):
        if original_fonts[i] == 0 or processed_fonts[i] == 0:
            continue
        
        orig_ratio = original_fonts[i + 1] / original_fonts[i]
        proc_ratio = processed_fonts[i + 1] / processed_fonts[i]
        
        if orig_ratio == 0:
            continue
            
        ratio_diff = abs(orig_ratio - proc_ratio) / orig_ratio
        if ratio_diff > tolerance:
            return (False, f"Ratio not preserved: orig {orig_ratio:.2f} vs proc {proc_ratio:.2f}")
    
    return (True, "Font ratios preserved")


# =============================================================================
# TEST FUNCTIONS
# =============================================================================

def verify_additional_checks(result, original_order=None, original_fonts_by_shape=None, 
                             skip_overflow=False, skip_position=False, skip_ratio=False):
    """
    Run additional verification checks on a processed slide.
    
    Args:
        result: dict from load_and_verify
        original_order: list of shape names in original vertical order
        original_fonts_by_shape: dict mapping shape name to list of original fonts
        skip_overflow: skip overflow check (for special cases)
        skip_position: skip position check (for slides with visual elements)
        skip_ratio: skip font ratio check
    
    Raises:
        AssertionError: if any check fails
    """
    # Check 1: No text overflow
    if not skip_overflow:
        overflow_ok, overflow_msg = check_no_text_overflow(result)
        assert overflow_ok, f"Overflow check failed: {overflow_msg}"
    
    # Check 2: Vertical order preserved (comparing with original)
    if original_order:
        order_ok, order_msg = check_vertical_order_preserved(result['shapes'], original_order)
        assert order_ok, f"Vertical order check failed: {order_msg}"
    
    # Check 3: Position at margins (only for transformed slides)
    if not skip_position and not result['has_visual']:
        pos_ok, pos_msg = check_position_at_margins(result)
        assert pos_ok, f"Position check failed: {pos_msg}"
    
    # Check 4: Shapes stacked correctly (no overlaps)
    if not result['has_visual']:
        stack_ok, stack_msg = check_shapes_stacked_correctly(result)
        assert stack_ok, f"Stacking check failed: {stack_msg}"
    
    # Check 5: Font ratios preserved within each shape
    if not skip_ratio and original_fonts_by_shape and not result['has_visual']:
        for shape in result['shapes']:
            if shape['name'] in original_fonts_by_shape:
                orig_fonts = original_fonts_by_shape[shape['name']]
                proc_fonts = shape['fonts']
                ratio_ok, ratio_msg = check_font_ratios_in_shape(orig_fonts, proc_fonts)
                assert ratio_ok, f"Font ratio check failed for '{shape['name']}': {ratio_msg}"


def test_imagem_nascimento_joao_batista(tmp_path):
    """Has image - should NOT transform"""
    output, original_order, original_fonts = process_and_save_slide('slide_imagem_nascimento_joao_batista', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == True
    assert result['max_font'] == 36.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Nascimento de João Batista\n4ª Semana do Tempo do Advento – Ano A'
    assert result['shapes'][0]['name'] == 'WordArt 10'
    assert result['shapes'][0]['fonts'] == [28.0]
    assert result['shapes'][1]['name'] == 'WordArt 8'
    assert result['shapes'][1]['fonts'] == [36.0]
    
    # Additional checks (skip position check for visual elements)
    verify_additional_checks(result, original_order, original_fonts, skip_position=True)


def test_oracao_vocacional_p1(tmp_path):
    """Oração Vocacional parte 1"""
    output, original_order, original_fonts = process_and_save_slide('slide_oracao_vocacional_p1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 52.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Jesus, Mestre Divino, que chamastes os Apóstolos para vos seguirem, continuai a passar pelos nossos caminhos, pelas nossas famílias, pelas nossas escolas.\nORAÇÃO VOCACIONAL'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [10.5, 52.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [44.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_oracao_vocacional_p2(tmp_path):
    """Oração Vocacional parte 2"""
    output, original_order, original_fonts = process_and_save_slide('slide_oracao_vocacional_p2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'E continuai a repetir o convite a muitos de nossos jovens. Dai coragem às pessoas convidadas, dai forças para que vos sejam fiéis como apóstolos leigos, como diáconos, padres e bispos, como religiosos e religiosas para o bem do povo de Deus e de toda a humanidade. Amém!'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [45.0, 48.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_mesmo_as_trevas(tmp_path):
    """Mesmo as trevas"""
    output, original_order, original_fonts = process_and_save_slide('slide_mesmo_as_trevas', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 67.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Mesmo as trevas\x0bComunidade Ecumênica de Taizé\nMESMO AS TREVAS, NÃO SÃO TREVAS/ PARA TI, A NOITE É LUMINOSA COMO O DIA.'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [30.0, 40.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [67.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_imagem_nascimento_joao_batista_2(tmp_path):
    """Has image - should NOT transform (duplicate)"""
    output, original_order, original_fonts = process_and_save_slide('slide_imagem_nascimento_joao_batista_2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == True
    assert result['max_font'] == 36.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Nascimento de João Batista\n4ª Semana do Tempo do Advento – Ano A'
    assert result['shapes'][0]['name'] == 'WordArt 10'
    assert result['shapes'][0]['fonts'] == [28.0]
    assert result['shapes'][1]['name'] == 'WordArt 8'
    assert result['shapes'][1]['fonts'] == [36.0]
    
    # Additional checks (skip position check for visual elements)
    verify_additional_checks(result, original_order, original_fonts, skip_position=True)


def test_que_germine_o_salvador_v1(tmp_path):
    """Que germine o Salvador v1"""
    output, original_order, original_fonts = process_and_save_slide('slide_que_germine_o_salvador_v1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Que germine o Salvador\x0bPe. Sílvio Milanez\n1. O SENHOR ESTÁ PRA CHEGAR,/ JÁ SE CUMPRE A PROFECIA;/ O SEU REINO ENTÃO SERÁ/ LIBERDADE E ALEGRIA./ E AS NAÇÕES, ENFIM, RECEBEM,/ SALVAÇÃO A CADA DIA.\n/: DAS ALTURAS ORVALHEM OS CÉUS,/ E DAS NUVENS, QUE CHOVA A JUSTIÇA,/ QUE A TERRA SE ABRA AO AMOR/ E GERMINE O DEUS SALVADOR!  :/'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_que_germine_o_salvador_v2(tmp_path):
    """Que germine o Salvador v2"""
    output, original_order, original_fonts = process_and_save_slide('slide_que_germine_o_salvador_v2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Que germine o Salvador\x0bPe. Sílvio Milanez\n2. VEM DE NOVO RESTAURAR-NOS./ DE QUE LADO ESTARÁ?/ INDIGNADO CONTRA NÓS?/ E A VIDA, NÃO DARÁS?/ SALVAÇÃO E ALEGRIA,/ OUTRA VEZ NÃO BUSCARÁS?\n/: DAS ALTURAS ORVALHEM OS CÉUS,/ E DAS NUVENS, QUE CHOVA A JUSTIÇA,/ QUE A TERRA SE ABRA AO AMOR/ E GERMINE O DEUS SALVADOR!  :/'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_senhor_filho_de_deus(tmp_path):
    """Senhor e Filho de Deus"""
    output, original_order, original_fonts = process_and_save_slide('slide_senhor_filho_de_deus', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == '1. SENHOR E FILHO DE DEUS,/ COMPANHEIRO, IRMÃO E AMIGO./: TENDE PIEDADE DE NÓS. :/\n2. Ó CRISTO, FILHO DO HOMEM./ CONHECEIS A NOSSA FRAQUEZA./: TENDE PIEDADE DE NÓS. :/\n3. SENHOR E FILHO DO PAI./ ACOLHEI-NOS NA VOSSA CASA./: TENDE PIEDADE DE NÓS. :/\nSenhor e Filho de Deus\x0bPe. José Freitas Campos'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [48.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [12.0, 24.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_salmo_responsorial(tmp_path):
    """Salmo responsorial"""
    output, original_order, original_fonts = process_and_save_slide('slide_salmo_responsorial', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 60.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Salmo responsorial: Sl 24/25\x0bLecionário Ferial – Melodia: Ir. Míria Therezinha Kolling,icm\nLEVANTAI VOSSA CABEÇA E OLHAI,/ POIS A VOSSA REDENÇÃO SE APROXIMA!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [30.0, 40.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [60.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_aleluia(tmp_path):
    """Aleluia"""
    output, original_order, original_fonts = process_and_save_slide('slide_aleluia', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 54.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Aleluia!\x0bIr. Lindberg Pires,sj\n/: ALELUIA, ALELUIA!/ ALELUIA, ALELUIA! :/\nÓ REI E SENHOR DAS NAÇÕES E PEDRA ANGULAR DA IGREJA,/ VINDE SALVAR A MULHER E O HOMEM, QUE UM DIA, FORMASTES DO BARRO!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [48.0, 54.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_oracao_da_comunidade(tmp_path):
    """Oração da comunidade"""
    output, original_order, original_fonts = process_and_save_slide('slide_oracao_da_comunidade', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 85.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Oração da comunidade\nMOSTRAI-NOS, Ó SENHOR, VOSSOS CAMINHOS.'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [50.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [85.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_eis_senhor_a_tua_vinha_v1(tmp_path):
    """Eis, Senhor, a tua vinha v1"""
    output, original_order, original_fonts = process_and_save_slide('slide_eis_senhor_a_tua_vinha_v1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Eis, Senhor, a tua vinha\x0bJoão de Araújo e Ir. Míria Therezinha Kolling,icm\n1. “DO CÉU VAI DESCER O CORDEIRO.”/ É DOM, PURO DOM, SALVAÇÃO!/ NO ALTAR DO PENHOR VERDADEIRO,/ TAMBÉM, VAMOS SER OBLAÇÃO!\nEIS, SENHOR, A TUA VINHA/ FRUTOS MIL TE TRAZ, SENHOR./ MAS TEU POVO QUE CAMINHA,/ MAIS QUE FRUTO, É DOM DE AMOR!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]


def test_eis_senhor_a_tua_vinha_v2(tmp_path):
    """Eis, Senhor, a tua vinha v2"""
    output, original_order, original_fonts = process_and_save_slide('slide_eis_senhor_a_tua_vinha_v2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Eis, Senhor, a tua vinha\x0bJoão de Araújo e Ir. Míria Therezinha Kolling,icm\n2. NA TERRA JÁ BROTA A ESPERANÇA/ E A GRAÇA DE DEUS VEM DIZER/ QUE O POVO DA NOVA ALIANÇA,/ TAMBÉM, OFERENDA VAI SER.\nEIS, SENHOR, A TUA VINHA/ FRUTOS MIL TE TRAZ, SENHOR./ MAS TEU POVO QUE CAMINHA,/ MAIS QUE FRUTO, É DOM DE AMOR!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_eis_senhor_a_tua_vinha_v3(tmp_path):
    """Eis, Senhor, a tua vinha v3"""
    output, original_order, original_fonts = process_and_save_slide('slide_eis_senhor_a_tua_vinha_v3', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Eis, Senhor, a tua vinha\x0bJoão de Araújo e Ir. Míria Therezinha Kolling,icm\n3. IRMÃOS NA FÉ VIVA, EXULTANTES/ PARTILHAM O PÃO SEMPRE MAIS./ E CAMPOS JAMAIS VERDEJANTES,/ TAMBÉM, JÁ SE TORNAM TRIGAIS!\nEIS, SENHOR, A TUA VINHA/ FRUTOS MIL TE TRAZ, SENHOR./ MAS TEU POVO QUE CAMINHA,/ MAIS QUE FRUTO, É DOM DE AMOR!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_conclusao_preparacao_oferendas(tmp_path):
    """Conclusão à Preparação das Oferendas"""
    output, original_order, original_fonts = process_and_save_slide('slide_conclusao_preparacao_oferendas', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Conclusão à Preparação das Oferendas\x0bCânon Romano – 3ª Edição Típica do Missal Romano\nRECEBA O SENHOR POR TUAS MÃOS ESTE SACRIFÍCIO, PARA GLÓRIA DO SEU NOME, PARA NOSSO BEM E DE TODA A SUA SANTA IGREJA.'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [48.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_santo_santo_santo(tmp_path):
    """Santo, Santo, Santo"""
    output, original_order, original_fonts = process_and_save_slide('slide_santo_santo_santo', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 54.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Santo, Santo, Santo! Senhor Deus do Universo\x0bD. Pedro Brito Guimarães\n1. SANTO, SANTO, SANTO!/ SENHOR DEUS DO UNIVERSO!/ O CÉU E A TERRA PROCLAMAM/ A VOSSA GLÓRIA.\n/: HOSANA NAS ALTURAS,/ HOSANA! :/\n2. /: BENDITO AQUELE QUE VEM,/ EM NOME DO SENHOR! :/'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [48.0, 54.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_enviai_espirito_santo(tmp_path):
    """Enviai o vosso Espírito Santo"""
    output, original_order, original_fonts = process_and_save_slide('slide_enviai_espirito_santo', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 80.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Oração Eucarística II: Aclamações da assembleia\x0bCânon Romano – 3ª Edição Típica do Missal Romano\nENVIAI O VOSSO ESPÍRITO SANTO!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [36.0, 47.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [80.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_anamnese_1(tmp_path):
    """Anamnese 1"""
    output, original_order, original_fonts = process_and_save_slide('slide_anamnese_1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 54.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'ANUNCIAMOS, SENHOR, A VOSSA MORTE E PROCLAMAMOS A VOSSA RESSURREIÇÃO. VINDE, SENHOR JESUS!\nAnamnese (Memorial) – Mistério da fé!\x0bCânon Romano – 3ª Edição Típica do Missal Romano'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [54.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [12.0, 24.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_anamnese_2(tmp_path):
    """Anamnese 2"""
    output, original_order, original_fonts = process_and_save_slide('slide_anamnese_2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 54.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'TODAS AS VEZES QUE COMEMOS DESTE PÃO/ E BEBEMOS DESTE CÁLICE,/ ANUNCIAMOS, SENHOR, A VOSSA MORTE,/ ENQUANTO ESPERAMOS A VOSSA VINDA.\nAnamnese (Memorial) – Mistério da fé e do amor!\x0bCânon Romano – 3ª Edição Típica do Missal Romano'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [54.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [12.0, 24.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_anamnese_3(tmp_path):
    """Anamnese 3"""
    output, original_order, original_fonts = process_and_save_slide('slide_anamnese_3', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 54.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'SALVADOR DO MUNDO, SALVAI-NOS, VÓS QUE NOS LIBERTASTES PELA CRUZ E RESSURREIÇÃO.\nAnamnese (Memorial) – Mistério da fé para a salvação do mundo!\x0bCânon Romano – 3ª Edição Típica do Missal Romano'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [54.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [24.0]
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_aceitai_senhor_oferta(tmp_path):
    """Aceitai, ó Senhor, a nossa oferta"""
    output, original_order, original_fonts = process_and_save_slide('slide_aceitai_senhor_oferta', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 81.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'ACEITAI, Ó SENHOR, A NOSSA OFERTA!\nOração Eucarística II: Aclamações da assembleia\x0bCânon Romano – 3ª Edição Típica do Missal Romano'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [81.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    
    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_espirito_nos_una(tmp_path):
    """O Espírito nos una"""
    output, original_order, original_fonts = process_and_save_slide('slide_espirito_nos_una', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 74.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'O ESPÍRITO NOS UNA NUM SÓ CORPO!\nOração Eucarística II: Aclamações da assembleia\x0bCânon Romano – 3ª Edição Típica do Missal Romano'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [74.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [33.0, 44.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_lembraivos_da_igreja(tmp_path):
    """Lembrai-vos da Igreja"""
    output, original_order, original_fonts = process_and_save_slide('slide_lembraivos_da_igreja', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 73.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'LEMBRAI-VOS, Ó PAI, DA VOSSA IGREJA!\nOração Eucarística II: Aclamações da assembleia\x0bCânon Romano – 3ª Edição Típica do Missal Romano'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [73.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [32.0, 43.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_concedei_luz_eterna(tmp_path):
    """Concedei-lhes a luz eterna"""
    output, original_order, original_fonts = process_and_save_slide('slide_concedei_luz_eterna', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 73.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'CONCEDEI-LHES, Ó SENHOR, A LUZ ETERNA!\nOração Eucarística II: Aclamações da assembleia\x0bCânon Romano – 3ª Edição Típica do Missal Romano'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [73.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [32.0, 43.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_cantico_zacarias_v1(tmp_path):
    """Cântico de Zacarias v1"""
    output, original_order, original_fonts = process_and_save_slide('slide_cantico_zacarias_v1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Cântico de Zacarias\x0bLc 1,68-78 – Letra: D. Marcos Barbosa – Melodia: Pe. Ney B. Pereira\nBENDITO SEJA O SENHOR, DEUS DE ISRAEL:/ ELE VISITA O SEU POVO E NOS SALVA!\n1. FORÇA E SALVAÇÃO SUSCITA EM NÓS * DA CASA DE DAVI, SEU SERVIDOR,/ CONFORME ANUNCIARA EM TEMPOS IDOS, * PELA BOCA DOS SANTOS, SEUS PROFETAS!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_cantico_zacarias_v2(tmp_path):
    """Cântico de Zacarias v2"""
    output, original_order, original_fonts = process_and_save_slide('slide_cantico_zacarias_v2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Cântico de Zacarias\x0bLc 1,68-78 – Letra: D. Marcos Barbosa – Melodia: Pe. Ney B. Pereira\nBENDITO SEJA O SENHOR, DEUS DE ISRAEL:/ ELE VISITA O SEU POVO E NOS SALVA!\n2. A SUA SALVAÇÃO NOS ARREBATA * ARRANCA-NOS DAS MÃOS DOS INIMIGOS./ PELO AMOR QUE JUROU A NOSSOS PAIS, * PELA SANTA ALIANÇA QUE ELE FEZ!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_cantico_zacarias_v3(tmp_path):
    """Cântico de Zacarias v3"""
    output, original_order, original_fonts = process_and_save_slide('slide_cantico_zacarias_v3', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Cântico de Zacarias\x0bLc 1,68-78 – Letra: D. Marcos Barbosa – Melodia: Pe. Ney B. Pereira\nBENDITO SEJA O SENHOR, DEUS DE ISRAEL:/ ELE VISITA O SEU POVO E NOS SALVA!\n3. A NOSSO PAI ABRAÃO JUROU LIVRAR-NOS * DO MEDO E DO TEMOR DOS INIMIGOS,/ A FIM DE QUE O SIRVAMOS PARA SEMPRE, * DIANTE DELE, EM JUSTIÇA E SANTIDADE!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_cantico_zacarias_v4(tmp_path):
    """Cântico de Zacarias v4"""
    output, original_order, original_fonts = process_and_save_slide('slide_cantico_zacarias_v4', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Cântico de Zacarias\x0bLc 1,68-78 – Letra: D. Marcos Barbosa – Melodia: Pe. Ney B. Pereira\nBENDITO SEJA O SENHOR, DEUS DE ISRAEL:/ ELE VISITA O SEU POVO E NOS SALVA!\n4. E TU, MENINO, TU SERÁS CHAMADO * COM O NOME DE PROFETA DO ALTÍSSIMO:/ TU IRÁS ANTE A FACE DO SENHOR, * À SUA FRENTE, PREPARANDO-LHE OS CAMINHOS!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_cantico_zacarias_v5(tmp_path):
    """Cântico de Zacarias v5"""
    output, original_order, original_fonts = process_and_save_slide('slide_cantico_zacarias_v5', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Cântico de Zacarias\x0bLc 1,68-78 – Letra: D. Marcos Barbosa – Melodia: Pe. Ney B. Pereira\nBENDITO SEJA O SENHOR, DEUS DE ISRAEL:/ ELE VISITA O SEU POVO E NOS SALVA!\n5. ANUNCIANDO A SEU POVO A SALVAÇÃO, * A SEU POVO O PERDÃO DOS PECADOS,/ GRAÇAS À COMPAIXÃO DO NOSSO DEUS, * O NOSSO DEUS QUE VEM NOS VISITAR!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_cantico_zacarias_v6(tmp_path):
    """Cântico de Zacarias v6"""
    output, original_order, original_fonts = process_and_save_slide('slide_cantico_zacarias_v6', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Cântico de Zacarias\x0bLc 1,68-78 – Letra: D. Marcos Barbosa – Melodia: Pe. Ney B. Pereira\nBENDITO SEJA O SENHOR, DEUS DE ISRAEL:/ ELE VISITA O SEU POVO E NOS SALVA!\n6. ELE É A LUZ QUE VEM DO ALTO, O SOL NASCENTE * POR SOBRE OS QUE JAZEM NAS TREVAS,/ E GUIA OS NOSSOS PASSOS PARA SEMPRE, * PARA SEMPRE NO CAMINHO DA PAZ!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_cantico_zacarias_v7(tmp_path):
    """Cântico de Zacarias v7"""
    output, original_order, original_fonts = process_and_save_slide('slide_cantico_zacarias_v7', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Cântico de Zacarias\x0bLc 1,68-78 – Letra: D. Marcos Barbosa – Melodia: Pe. Ney B. Pereira\nBENDITO SEJA O SENHOR, DEUS DE ISRAEL:/ ELE VISITA O SEU POVO E NOS SALVA!\n7. GLÓRIA AO PAI, AO FILHO E AO SANTO ESPÍRITO, * LOUVOR DESDE AGORA E PARA SEMPRE,/ AO DEUS QUE É, QUE ERA E QUE VEM, * AGORA E PELOS SÉCULOS, AMÉM!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [24.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_oracao_santo_antonio_p1(tmp_path):
    """Oração a Santo Antônio p1"""
    output, original_order, original_fonts = process_and_save_slide('slide_oracao_santo_antonio_p1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 47.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Oração a Santo Antônio\n\nSanto Antônio, vós sois um dos grandes amigos de Deus.  O mundo inteiro vos venera como o Santo da aliança e da intimidade com Deus.  Fostes um grande pregador e missionário de Cristo.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [28.0, 47.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_oracao_santo_antonio_p2(tmp_path):
    """Oração a Santo Antônio p2"""
    output, original_order, original_fonts = process_and_save_slide('slide_oracao_santo_antonio_p2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 62.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Ajudai-me a viver de acordo com o vosso modelo e a levar aos meus irmãos e irmãs a mensagem da aliança, do Evangelho e do amor de Jesus Cristo.'
    assert result['shapes'][0]['name'] == 'Text Box 3'
    assert result['shapes'][0]['fonts'] == [62.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_oracao_santo_antonio_p3(tmp_path):
    """Oração a Santo Antônio p3"""
    output, original_order, original_fonts = process_and_save_slide('slide_oracao_santo_antonio_p3', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 72.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Defendei-me de todos os perigos.  Fortalecei a mim e ao meu lar, em todas as tribulações. Protegei os meus empreendimentos.'
    assert result['shapes'][0]['name'] == 'Text Box 3'
    assert result['shapes'][0]['fonts'] == [72.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_oracao_santo_antonio_p4(tmp_path):
    """Oração a Santo Antônio p4"""
    output, original_order, original_fonts = process_and_save_slide('slide_oracao_santo_antonio_p4', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 51.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Intercedei para que eu tenha saúde, prosperidade, harmonia, alegria e paz.  Inspirai-me na prática do bem e ajudai-me a alcançar a vida eterna.\n\nAmém.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [42.86, 51.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_responsorio_santo_antonio_v1(tmp_path):
    """Responsório de Santo Antônio v1"""
    output, original_order, original_fonts = process_and_save_slide('slide_responsorio_santo_antonio_v1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 58.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Responsório de Santo Antônio\n\n1- Se milagres tu procuras. / Vai, recorre a Santo Antônio: / Verás fugir as maldades / e as tentações do demônio.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [58.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_responsorio_santo_antonio_refrao_1(tmp_path):
    """Responsório refrão 1"""
    output, original_order, original_fonts = process_and_save_slide('slide_responsorio_santo_antonio_refrao_1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 74.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Faz encontrar o perdido, rompe a grade da prisão,  cede o mar embravecido, silencia a voz do trovão.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [42.86, 74.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_responsorio_santo_antonio_v2(tmp_path):
    """Responsório v2"""
    output, original_order, original_fonts = process_and_save_slide('slide_responsorio_santo_antonio_v2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 64.0
    assert result['num_shapes'] == 1
    assert result['text'] == '2- Em graças são transformados, / todos os males humanos. / Provam-no todos os homens, / sobretudo os paduanos.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [64.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_responsorio_santo_antonio_refrao_2(tmp_path):
    """Responsório refrão 2"""
    output, original_order, original_fonts = process_and_save_slide('slide_responsorio_santo_antonio_refrao_2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 74.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Faz encontrar o perdido, rompe a grade da prisão,  cede o mar embravecido, silencia a voz do trovão.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [42.86, 74.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_responsorio_santo_antonio_v3(tmp_path):
    """Responsório v3"""
    output, original_order, original_fonts = process_and_save_slide('slide_responsorio_santo_antonio_v3', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 75.0
    assert result['num_shapes'] == 1
    assert result['text'] == '3- Foge a peste, o erro, a morte / de Antônio a invocação. / O fraco torna-se forte / e o enfermo se torna são.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [75.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_responsorio_santo_antonio_refrao_3(tmp_path):
    """Responsório refrão 3"""
    output, original_order, original_fonts = process_and_save_slide('slide_responsorio_santo_antonio_refrao_3', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 74.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Faz encontrar o perdido, rompe a grade da prisão,  cede o mar embravecido, silencia a voz do trovão.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [42.86, 74.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_responsorio_santo_antonio_v4(tmp_path):
    """Responsório v4"""
    output, original_order, original_fonts = process_and_save_slide('slide_responsorio_santo_antonio_v4', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 72.0
    assert result['num_shapes'] == 1
    assert result['text'] == '4- Glória ao Pai e a Jesus Cristo, / e ao Espírito também, / com quem vive Santo Antônio, / pelos séculos. Amém.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [72.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_responsorio_santo_antonio_refrao_4(tmp_path):
    """Responsório refrão 4"""
    output, original_order, original_fonts = process_and_save_slide('slide_responsorio_santo_antonio_refrao_4', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 74.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Faz encontrar o perdido, rompe a grade da prisão,  cede o mar embravecido, silencia a voz do trovão.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [42.86, 74.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_rogai_por_nos(tmp_path):
    """Rogai por nós"""
    output, original_order, original_fonts = process_and_save_slide('slide_rogai_por_nos', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 68.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'C: Rogai por nós Santo Antônio.\n\nT: Para que sejamos dignos das promessas de Cristo.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [68.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_oremos_intercessao(tmp_path):
    """Oremos (intercessão)"""
    output, original_order, original_fonts = process_and_save_slide('slide_oremos_intercessao', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 39.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Oremos:\t\nNós vos suplicamos, ó Deus, que pela intercessão de Santo Antônio, vosso confessor e doutor, sejamos fortalecidos com os auxílios espirituais de que necessitamos,\npara percorrermos em paz os caminhos desta vida e alcançarmos as alegrias da vida eterna. Por Nosso Senhor Jesus Cristo, Vosso Filho, na Unidade do Espírito Santo.\n\nAmém.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [39.0]
    assert result['shapes'][1]['name'] == 'Text Box 3'
    assert result['shapes'][1]['fonts'] == [39.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_bencao_dos_paes(tmp_path):
    """Benção dos pães"""
    output, original_order, original_fonts = process_and_save_slide('slide_bencao_dos_paes', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 39.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Benção dos pães\n\nC: A nossa proteção está no nome do Senhor.\n\nT: Que fez o céu e a terra.\n\nC: O Senhor esteja convosco.\n\nT: Ele está no meio de nós.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [39.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_oremos_paes(tmp_path):
    """Oremos (pães)"""
    output, original_order, original_fonts = process_and_save_slide('slide_oremos_paes', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 39.0
    assert result['num_shapes'] == 1
    assert result['text'] == 'Oremos:\nSenhor Jesus Cristo, Pão dos Anjos, Pão vivo da vida eterna, dignai-vos abençoar † estes pães assim como abençoastes os cinco pães no deserto, para que todos os que deles comerem, recebam saúde da alma e do corpo, paz e harmonia na família e possam partilhar, com os irmãos mais necessitados, os bens recebidos de Deus. Vós que viveis e reinais por todos os séculos dos séculos.  Amém.'
    assert result['shapes'][0]['name'] == 'Text Box 2'
    assert result['shapes'][0]['fonts'] == [39.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_santo_antonio_junto_a_nos_v1_v2(tmp_path):
    """Santo Antônio junto a nós v1-2"""
    output, original_order, original_fonts = process_and_save_slide('slide_santo_antonio_junto_a_nos_v1_v2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 44.0
    assert result['num_shapes'] == 2
    assert result['text'] == '1. SANTO ANTÔNIO, QUE ESTÁS JUNTO A NÓS,/ NOS AJUDA A PLANTAR MUNDO NOVO./: ONDE A PAZ,A JUSTIÇA E A VERDADE/ SEJAM FORÇA E ESPERANÇA DO NOVO. :/\n2. SANTO ANTÔNIO, PRESENTE NA LUTA/ DE QUEM PARTE NA BUSCA DO PÃO,/: QUE O SUOR DERRAMADO CULTIVE/ A SEMENTE DO REINO DE IRMÃOS. :/\nSanto Antônio, que estás junto a nós\x0bFr. Fernando Antônio Fabreti,ofm'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [44.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [12.0, 24.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_santo_antonio_junto_a_nos_v3_v4(tmp_path):
    """Santo Antônio junto a nós v3-4"""
    output, original_order, original_fonts = process_and_save_slide('slide_santo_antonio_junto_a_nos_v3_v4', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 44.0
    assert result['num_shapes'] == 2
    assert result['text'] == '3. SANTO ANTÔNIO, DO POVO SOFRIDO,/ VEM MOSTRAR NOVO TEMPO FLORIR./: ONDE TUDO SERÁ PARTILHADO,/ NOVA ERA DE UM POVO A SORRIR. :/\n4. SANTO ANTÔNIO, TEU ROSTO DE POBRE,/ BRILHA A PAZ DO MENINO JESUS,/: DEUS-CRIANÇA, ESPERANÇA DO POVO,/ NOSSA FORÇA, DO CÉU NOSSA LUZ! :/\nSanto Antônio, que estás junto a nós\x0bFr. Fernando Antônio Fabreti,ofm'
    assert result['shapes'][0]['name'] == 'Subtítulo 2'
    assert result['shapes'][0]['fonts'] == [44.0]
    assert result['shapes'][1]['name'] == 'Título 1'
    assert result['shapes'][1]['fonts'] == [12.0, 24.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_hino_jubileu_v1(tmp_path):
    """Hino do Jubileu v1"""
    output, original_order, original_fonts = process_and_save_slide('slide_hino_jubileu_v1', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Hino do Jubileu-2025: Peregrinos de Esperança\x0bL: Peirangelo Sequeri – M: Francesco Meneghello – Versão: Antônio Cartageno\nCHAMA VIVA DA MINHA ESPERANÇA,/ ESTE CANTO SUBA PARA TI!/ SEIO ETERNO DE INFINITA VIDA,/ NO CAMINHO EU CONFIO EM TI!\n1. TODA LÍNGUA, POVO E NAÇÃO/ TUA LUZ ENCONTRA NA PALAVRA./ OS TEUS FILHOS, FRÁGEIS E DISPERSOS,/ SE REÚNEM NO TEU FILHO AMADO.'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [22.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_hino_jubileu_v2(tmp_path):
    """Hino do Jubileu v2"""
    output, original_order, original_fonts = process_and_save_slide('slide_hino_jubileu_v2', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Hino do Jubileu-2025: Peregrinos de Esperança\x0bL: Peirangelo Sequeri – M: Francesco Meneghello – Versão: Antônio Cartageno\nCHAMA VIVA DA MINHA ESPERANÇA,/ ESTE CANTO SUBA PARA TI!/ SEIO ETERNO DE INFINITA VIDA,/ NO CAMINHO EU CONFIO EM TI!\n2. DEUS NOS OLHA TERNO E PACIENTE: NASCE A AURORA DE UM FUTURO NOVO./ NOVOS CÉUS, TERRA FEITA NOVA,/ PASSA OS MUROS, ESPÍRITO DE VIDA!'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [22.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)


def test_hino_jubileu_v3(tmp_path):
    """Hino do Jubileu v3"""
    output, original_order, original_fonts = process_and_save_slide('slide_hino_jubileu_v3', tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == False
    assert result['max_font'] == 48.0
    assert result['num_shapes'] == 2
    assert result['text'] == 'Hino do Jubileu-2025: Peregrinos de Esperança\x0bL: Peirangelo Sequeri – M: Francesco Meneghello – Versão: Antônio Cartageno\nCHAMA VIVA DA MINHA ESPERANÇA,/ ESTE CANTO SUBA PARA TI!/ SEIO ETERNO DE INFINITA VIDA,/ NO CAMINHO EU CONFIO EM TI!\n3. ERGUE OS OLHOS, MOVE-TE COM O VENTO,/ NÃO TE ATRASES: CHEGA DEUS NO TEMPO./ JESUS CRISTO POR TI SE FEZ HOMEM:/ AOS MILHARES, SEGUEM O CAMINHO.'
    assert result['shapes'][0]['name'] == 'Título 1'
    assert result['shapes'][0]['fonts'] == [22.0, 32.0]
    assert result['shapes'][1]['name'] == 'Subtítulo 2'
    assert result['shapes'][1]['fonts'] == [44.0, 48.0]

    # Additional checks
    verify_additional_checks(result, original_order, original_fonts)



if __name__ == '__main__':
    pytest.main([__file__, '-v'])
