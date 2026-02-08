"""
Unit tests for fix_slides_for_obs_processor.

These tests verify the processor's utility functions without requiring
external presentation files. For slide-specific tests, see test_individual_slides.py.

Run with: python -m pytest tests/test_slides_processor.py -v
Or from tests/: python -m pytest test_slides_processor.py -v
"""
import unittest
import os
import sys

# Add parent directory to path for imports
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import fix_slides_for_obs_processor as processor


class TestTextMeasurement(unittest.TestCase):
    """Test text measurement functions."""
    
    def test_measure_single_line(self):
        """Test measuring a single line of text."""
        text = "Hello World"
        size = processor.measure_multiline_text_size(text, "Arial", 48)
        self.assertIsNotNone(size)
        self.assertGreater(size[0], 0)  # Width > 0
        self.assertGreater(size[1], 0)  # Height > 0
    
    def test_measure_multiline(self):
        """Test measuring multi-line text."""
        text = "Line 1\nLine 2\nLine 3"
        size = processor.measure_multiline_text_size(text, "Arial", 48)
        self.assertIsNotNone(size)
        
        # Multi-line should be taller than single line
        single_size = processor.measure_multiline_text_size("Line 1", "Arial", 48)
        self.assertGreater(size[1], single_size[1])
    
    def test_measure_with_wrapping(self):
        """Test text wrapping to max width."""
        text = "This is a very long line that should wrap to multiple lines when constrained"
        max_width = 200
        size = processor.measure_multiline_text_size(text, "Arial", 24, max_width)
        self.assertIsNotNone(size)
        # Width should not exceed max_width (with some tolerance for word boundaries)
        self.assertLessEqual(size[0], max_width + 50)
    
    def test_normalize_whitespace(self):
        """Test text whitespace normalization."""
        text_with_blanks = "Line 1\n\n\n\nLine 2"
        normalized = processor.normalize_text_whitespace(text_with_blanks)
        # Should reduce multiple blank lines to at most one
        self.assertNotIn("\n\n\n", normalized)
    
    def test_measure_empty_text(self):
        """Test measuring empty text returns valid result."""
        size = processor.measure_multiline_text_size("", "Arial", 48)
        # Empty text should return some result (height for empty line)
        self.assertIsNotNone(size)
    
    def test_measure_different_fonts(self):
        """Test measuring with different font sizes."""
        text = "Test text"
        size_small = processor.measure_multiline_text_size(text, "Arial", 12)
        size_large = processor.measure_multiline_text_size(text, "Arial", 48)
        
        self.assertIsNotNone(size_small)
        self.assertIsNotNone(size_large)
        # Larger font should produce larger dimensions
        self.assertGreater(size_large[0], size_small[0])
        self.assertGreater(size_large[1], size_small[1])


class TestFontPath(unittest.TestCase):
    """Test font path resolution."""
    
    def test_arial_font_path(self):
        """Test that Arial font can be resolved."""
        path = processor.get_font_path("Arial")
        self.assertIsNotNone(path)
        self.assertTrue(os.path.exists(path) or path.endswith('.ttf') or path.endswith('.ttc'))
    
    def test_unknown_font_fallback(self):
        """Test that unknown fonts fall back to a default."""
        path = processor.get_font_path("NonExistentFont12345")
        # Should return some path (fallback to Arial or default)
        self.assertIsNotNone(path)


class TestTextNormalization(unittest.TestCase):
    """Test text normalization functions."""
    
    def test_normalize_multiple_newlines(self):
        """Multiple consecutive newlines should be reduced."""
        text = "Line 1\n\n\n\nLine 2"
        normalized = processor.normalize_text_whitespace(text)
        # Should not have more than 2 consecutive newlines
        self.assertNotIn("\n\n\n", normalized)
    
    def test_normalize_preserves_single_newlines(self):
        """Single newlines should be preserved."""
        text = "Line 1\nLine 2\nLine 3"
        normalized = processor.normalize_text_whitespace(text)
        self.assertEqual(normalized.count('\n'), 2)
    
    def test_normalize_strips_whitespace(self):
        """Leading and trailing whitespace should be stripped."""
        text = "  \n  Text content  \n  "
        normalized = processor.normalize_text_whitespace(text)
        self.assertFalse(normalized.startswith(' '))
        self.assertFalse(normalized.endswith(' '))


class TestEMUConversions(unittest.TestCase):
    """Test EMU conversion constants and calculations."""
    
    def test_emu_per_point(self):
        """Test EMU to point conversion."""
        # 1 point = 12700 EMUs
        emu_per_pt = 12700
        
        # 72 points = 1 inch = 914400 EMUs
        points = 72
        expected_emu = points * emu_per_pt
        self.assertEqual(expected_emu, 914400)
    
    def test_standard_slide_dimensions(self):
        """Test standard 16:9 slide dimensions."""
        # Standard 16:9 slide is 10" x 5.625" (or 960pt x 540pt)
        slide_width_pt = 960
        slide_height_pt = 540
        
        # Aspect ratio should be 16:9
        ratio = slide_width_pt / slide_height_pt
        expected_ratio = 16 / 9
        self.assertAlmostEqual(ratio, expected_ratio, places=2)


class MockPlaceholderFormat:
    """Mock placeholder format for testing."""
    def __init__(self, placeholder_type):
        self.type = placeholder_type


class MockShape:
    """Mock shape for testing is_insignificant_placeholder."""
    def __init__(self, is_placeholder=False, placeholder_type=None):
        self._is_placeholder = is_placeholder
        self._placeholder_type = placeholder_type
    
    @property
    def is_placeholder(self):
        return self._is_placeholder
    
    @property
    def placeholder_format(self):
        if self._placeholder_type is not None:
            return MockPlaceholderFormat(self._placeholder_type)
        return None


class TestInsignificantPlaceholder(unittest.TestCase):
    """Test is_insignificant_placeholder function."""
    
    def test_non_placeholder_shape(self):
        """Non-placeholder shapes should return False."""
        shape = MockShape(is_placeholder=False)
        self.assertFalse(processor.is_insignificant_placeholder(shape))
    
    def test_slide_number_placeholder(self):
        """Slide number placeholders should be insignificant."""
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MockShape(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.SLIDE_NUMBER)
        self.assertTrue(processor.is_insignificant_placeholder(shape))
    
    def test_footer_placeholder(self):
        """Footer placeholders should be insignificant."""
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MockShape(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.FOOTER)
        self.assertTrue(processor.is_insignificant_placeholder(shape))
    
    def test_date_placeholder(self):
        """Date placeholders should be insignificant."""
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MockShape(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.DATE)
        self.assertTrue(processor.is_insignificant_placeholder(shape))
    
    def test_header_placeholder(self):
        """Header placeholders should be insignificant."""
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MockShape(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.HEADER)
        self.assertTrue(processor.is_insignificant_placeholder(shape))
    
    def test_title_placeholder_is_significant(self):
        """Title placeholders should be significant (return False)."""
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MockShape(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.TITLE)
        self.assertFalse(processor.is_insignificant_placeholder(shape))
    
    def test_body_placeholder_is_significant(self):
        """Body placeholders should be significant (return False)."""
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MockShape(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.BODY)
        self.assertFalse(processor.is_insignificant_placeholder(shape))
    
    def test_subtitle_placeholder_is_significant(self):
        """Subtitle placeholders should be significant (return False)."""
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MockShape(is_placeholder=True, placeholder_type=PP_PLACEHOLDER.SUBTITLE)
        self.assertFalse(processor.is_insignificant_placeholder(shape))
    
    def test_placeholder_without_format(self):
        """Placeholder without format should return False (not crash)."""
        shape = MockShape(is_placeholder=True, placeholder_type=None)
        # Should return False, not raise an exception
        self.assertFalse(processor.is_insignificant_placeholder(shape))


class TestColorInversion(unittest.TestCase):
    """Test color inversion logic used when invert_colors option is enabled."""
    
    def invert_color(self, color_hex):
        """Helper to invert a hex color string (mimics processor logic)."""
        color_hex = color_hex.lstrip('#')
        r = int(color_hex[0:2], 16)
        g = int(color_hex[2:4], 16)
        b = int(color_hex[4:6], 16)
        return f"{255-r:02X}{255-g:02X}{255-b:02X}"
    
    def test_invert_black_to_white(self):
        """Inverting black (#000000) should give white (#FFFFFF)."""
        result = self.invert_color("#000000")
        self.assertEqual(result, "FFFFFF")
    
    def test_invert_white_to_black(self):
        """Inverting white (#FFFFFF) should give black (#000000)."""
        result = self.invert_color("#FFFFFF")
        self.assertEqual(result, "000000")
    
    def test_invert_almost_black(self):
        """Inverting almost black (#010101) should give almost white (#FEFEFE)."""
        result = self.invert_color("#010101")
        self.assertEqual(result, "FEFEFE")
    
    def test_invert_almost_white(self):
        """Inverting almost white (#FEFEFE) should give almost black (#010101)."""
        result = self.invert_color("#FEFEFE")
        self.assertEqual(result, "010101")
    
    def test_invert_default_glow_color(self):
        """Inverting default glow color (#FFFFF0) should give (#00000F)."""
        result = self.invert_color("#FFFFF0")
        self.assertEqual(result, "00000F")
    
    def test_invert_red(self):
        """Inverting pure red (#FF0000) should give cyan (#00FFFF)."""
        result = self.invert_color("#FF0000")
        self.assertEqual(result, "00FFFF")
    
    def test_invert_green(self):
        """Inverting pure green (#00FF00) should give magenta (#FF00FF)."""
        result = self.invert_color("#00FF00")
        self.assertEqual(result, "FF00FF")
    
    def test_invert_blue(self):
        """Inverting pure blue (#0000FF) should give yellow (#FFFF00)."""
        result = self.invert_color("#0000FF")
        self.assertEqual(result, "FFFF00")
    
    def test_invert_gray(self):
        """Inverting mid-gray (#808080) should give complementary gray (#7F7F7F)."""
        result = self.invert_color("#808080")
        self.assertEqual(result, "7F7F7F")
    
    def test_invert_without_hash(self):
        """Inversion should work with or without # prefix."""
        result = self.invert_color("FF0000")
        self.assertEqual(result, "00FFFF")
    
    def test_invert_is_reversible(self):
        """Inverting twice should return to original color."""
        original = "A5B7C9"
        inverted = self.invert_color(original)
        double_inverted = self.invert_color(inverted)
        self.assertEqual(double_inverted, original)


class TestProcessPresentationInvertColors(unittest.TestCase):
    """Test process_presentation with invert_colors option."""
    
    def setUp(self):
        """Create a minimal test presentation."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        self.prs = Presentation()
        # Set standard slide dimensions
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
        # Add a slide with text
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        self.slide = self.prs.slides.add_slide(blank_layout)
        
        # Add a text box
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(2)
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Test text"
        run.font.size = Pt(24)
    
    def test_process_without_invert(self):
        """Process presentation without invert_colors should use white background."""
        from pptx.dml.color import RGBColor
        
        processor.process_presentation(
            self.prs, 
            glow_color="#FFFFF0", 
            glow_size=20, 
            text_color="010101",
            invert_colors=False
        )
        
        # Check background is white
        fill = self.slide.background.fill
        self.assertEqual(fill.fore_color.rgb, RGBColor(255, 255, 255))
    
    def test_process_with_invert(self):
        """Process presentation with invert_colors should use black background."""
        from pptx.dml.color import RGBColor
        
        processor.process_presentation(
            self.prs, 
            glow_color="#FFFFF0", 
            glow_size=20, 
            text_color="010101",
            invert_colors=True
        )
        
        # Check background is black
        fill = self.slide.background.fill
        self.assertEqual(fill.fore_color.rgb, RGBColor(0, 0, 0))
    
    def test_text_color_normal_mode(self):
        """Text color should match input in normal mode."""
        from pptx.dml.color import RGBColor
        
        processor.process_presentation(
            self.prs, 
            glow_color="#FFFFF0", 
            glow_size=20, 
            text_color="010101",
            invert_colors=False
        )
        
        # Find text and check color
        for shape in self.slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            # Should be almost black (#010101)
                            self.assertEqual(run.font.color.rgb, RGBColor(1, 1, 1))
    
    def test_text_color_inverted_mode(self):
        """Text color should be inverted in inverted mode."""
        from pptx.dml.color import RGBColor
        
        processor.process_presentation(
            self.prs, 
            glow_color="#FFFFF0", 
            glow_size=20, 
            text_color="010101",
            invert_colors=True
        )
        
        # Find text and check color
        for shape in self.slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            # Should be almost white (#FEFEFE = 254, 254, 254)
                            self.assertEqual(run.font.color.rgb, RGBColor(254, 254, 254))
    
    def test_empty_slide_background_normal(self):
        """Empty slides should have black background in normal mode."""
        from pptx.dml.color import RGBColor
        
        # Add an empty slide
        blank_layout = self.prs.slide_layouts[6]
        empty_slide = self.prs.slides.add_slide(blank_layout)
        
        processor.process_presentation(
            self.prs, 
            glow_color="#FFFFF0", 
            glow_size=20, 
            text_color="010101",
            invert_colors=False
        )
        
        # Empty slide should have black background
        fill = empty_slide.background.fill
        self.assertEqual(fill.fore_color.rgb, RGBColor(0, 0, 0))
    
    def test_empty_slide_background_inverted(self):
        """Empty slides should have white background in inverted mode."""
        from pptx.dml.color import RGBColor
        
        # Add an empty slide
        blank_layout = self.prs.slide_layouts[6]
        empty_slide = self.prs.slides.add_slide(blank_layout)
        
        processor.process_presentation(
            self.prs, 
            glow_color="#FFFFF0", 
            glow_size=20, 
            text_color="010101",
            invert_colors=True
        )
        
        # Empty slide should have white background in inverted mode
        fill = empty_slide.background.fill
        self.assertEqual(fill.fore_color.rgb, RGBColor(255, 255, 255))
    
    def test_default_invert_colors_is_false(self):
        """Default value for invert_colors should be False."""
        from pptx.dml.color import RGBColor
        
        # Call without invert_colors parameter
        processor.process_presentation(
            self.prs, 
            glow_color="#FFFFF0", 
            glow_size=20, 
            text_color="010101"
        )
        
        # Should behave as if invert_colors=False (white background)
        fill = self.slide.background.fill
        self.assertEqual(fill.fore_color.rgb, RGBColor(255, 255, 255))


if __name__ == '__main__':
    unittest.main(verbosity=2)
