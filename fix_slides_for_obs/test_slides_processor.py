"""
Unit tests for fix_slides_for_obs_processor.

Tests verify that each slide is either:
1. Transformed correctly (text repositioned and font scaled)
2. Skipped (if it contains visual elements like images/charts)
3. Unchanged (if no text shapes to process)

Run with: python -m pytest test_slides_processor.py -v
Or: python test_slides_processor.py
"""
import unittest
import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fix_slides_for_obs_processor as processor


class TestSlidesProcessor(unittest.TestCase):
    """Test suite for slide processing."""
    
    @classmethod
    def setUpClass(cls):
        """Load presentations once for all tests."""
        cls.prs_orig = Presentation('Apresentação1Original.pptx')
        cls.prs_proc = Presentation('test_output.pptx')
        cls.slides_orig = list(cls.prs_orig.slides)
        cls.slides_proc = list(cls.prs_proc.slides)
        
        # Load expected data
        with open('test_expected_data.json', 'r', encoding='utf-8') as f:
            cls.expected_data = json.load(f)
        
        # Slide dimensions
        cls.slide_width = cls.prs_orig.slide_width
        cls.slide_height = cls.prs_orig.slide_height
        cls.margin_x = int(cls.slide_width * 0.05)
        cls.margin_y = int(cls.slide_height * 0.05)
        cls.available_width = cls.slide_width - (2 * cls.margin_x)
        cls.available_height = cls.slide_height - (2 * cls.margin_y)
    
    def get_slide_text_shapes(self, slide):
        """Get text shapes from a slide."""
        shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                shapes.append(shape)
        return shapes
    
    def get_shape_font_size(self, shape):
        """Get the first font size found in a shape."""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() and run.font.size:
                    return run.font.size.pt
        return None
    
    def slide_has_visual_elements(self, slide):
        """Check if slide has visual elements."""
        visual_types = {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.MEDIA,
            MSO_SHAPE_TYPE.TABLE,
            MSO_SHAPE_TYPE.GROUP,
            MSO_SHAPE_TYPE.DIAGRAM,
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
            MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
        }
        for shape in slide.shapes:
            if shape.shape_type in visual_types:
                return True
        return False


# Generate individual test methods for each slide
def generate_slide_test(slide_num, expected):
    """Generate a test method for a specific slide."""
    
    def test_slide(self):
        slide_orig = self.slides_orig[slide_num - 1]
        slide_proc = self.slides_proc[slide_num - 1]
        
        orig_shapes = self.get_slide_text_shapes(slide_orig)
        proc_shapes = self.get_slide_text_shapes(slide_proc)
        
        # Test 1: Visual elements detection
        has_visual = self.slide_has_visual_elements(slide_orig)
        self.assertEqual(
            has_visual, 
            expected['has_visual_elements'],
            f"Slide {slide_num}: Visual elements detection mismatch"
        )
        
        # Test 2: Number of text shapes preserved
        self.assertEqual(
            len(proc_shapes),
            expected['num_text_shapes_proc'],
            f"Slide {slide_num}: Number of text shapes changed unexpectedly"
        )
        
        if expected['was_transformed']:
            # Test 3: For transformed slides, verify at least one shape is at margin position
            if proc_shapes:
                # Sort shapes by top position to get the topmost one
                shapes_sorted = sorted(proc_shapes, key=lambda s: s.top)
                top_shape = shapes_sorted[0]
                
                # The topmost shape should be at or near margin_y
                self.assertAlmostEqual(
                    top_shape.left, 
                    self.margin_x, 
                    delta=1000,  # Allow small variance
                    msg=f"Slide {slide_num}: Left position not at margin"
                )
                self.assertAlmostEqual(
                    top_shape.top, 
                    self.margin_y, 
                    delta=1000,
                    msg=f"Slide {slide_num}: Top position not at margin"
                )
                
                # Test 4: Width should match available width for the topmost shape
                self.assertAlmostEqual(
                    top_shape.width, 
                    self.available_width, 
                    delta=1000,
                    msg=f"Slide {slide_num}: Width not matching available width"
                )
                
                # Test 5: Font size should match expected (check any processed shape)
                if 'proc_font_pt' in expected:
                    # Find the shape with the expected font size
                    found_font = False
                    for shape in proc_shapes:
                        font_size = self.get_shape_font_size(shape)
                        if font_size and abs(font_size - expected['proc_font_pt']) < 2.0:
                            found_font = True
                            break
                    self.assertTrue(
                        found_font,
                        f"Slide {slide_num}: Expected font size {expected['proc_font_pt']}pt not found"
                    )
        
        elif expected['has_visual_elements']:
            # Test 6: Slides with visual elements should be unchanged
            if orig_shapes and proc_shapes:
                orig_shape = orig_shapes[0]
                proc_shape = proc_shapes[0]
                
                # Position should be unchanged
                self.assertEqual(
                    orig_shape.left, 
                    proc_shape.left,
                    f"Slide {slide_num}: Should be skipped but position changed"
                )
                self.assertEqual(
                    orig_shape.top, 
                    proc_shape.top,
                    f"Slide {slide_num}: Should be skipped but position changed"
                )
                
                # Font size should be unchanged
                orig_font = self.get_shape_font_size(orig_shape)
                proc_font = self.get_shape_font_size(proc_shape)
                if orig_font and proc_font:
                    self.assertAlmostEqual(
                        orig_font,
                        proc_font,
                        delta=0.5,
                        msg=f"Slide {slide_num}: Should be skipped but font changed"
                    )
    
    return test_slide


# Dynamically add test methods for each slide
def load_tests():
    """Load expected data and generate test methods."""
    if os.path.exists('test_expected_data.json'):
        with open('test_expected_data.json', 'r', encoding='utf-8') as f:
            expected_data = json.load(f)
        
        for expected in expected_data:
            slide_num = expected['slide_num']
            test_name = f'test_slide_{slide_num:02d}'
            test_method = generate_slide_test(slide_num, expected)
            test_method.__doc__ = f"Test slide {slide_num}"
            setattr(TestSlidesProcessor, test_name, test_method)


# Load tests before running
load_tests()


class TestTextMeasurement(unittest.TestCase):
    """Test text measurement functions."""
    
    def test_measure_single_line(self):
        """Test measuring a single line of text."""
        text = "Hello World"
        size = processor.measure_multiline_text_size(text, "Arial", 48)
        self.assertIsNotNone(size)
        self.assertGreater(size[0], 0)  # Width > 0
        self.assertGreater(size[1], 0)  # Height > 0
    
    def test_measure_multiline(self):
        """Test measuring multi-line text."""
        text = "Line 1\nLine 2\nLine 3"
        size = processor.measure_multiline_text_size(text, "Arial", 48)
        self.assertIsNotNone(size)
        
        # Multi-line should be taller than single line
        single_size = processor.measure_multiline_text_size("Line 1", "Arial", 48)
        self.assertGreater(size[1], single_size[1])
    
    def test_measure_with_wrapping(self):
        """Test text wrapping to max width."""
        text = "This is a very long line that should wrap to multiple lines when constrained"
        max_width = 200
        size = processor.measure_multiline_text_size(text, "Arial", 24, max_width)
        self.assertIsNotNone(size)
        # Width should not exceed max_width (with some tolerance for word boundaries)
        self.assertLessEqual(size[0], max_width + 50)
    
    def test_normalize_whitespace(self):
        """Test text whitespace normalization."""
        text_with_blanks = "Line 1\n\n\n\nLine 2"
        normalized = processor.normalize_text_whitespace(text_with_blanks)
        # Should reduce multiple blank lines to at most one
        self.assertNotIn("\n\n\n", normalized)


class TestVisualElementDetection(unittest.TestCase):
    """Test visual element detection."""
    
    @classmethod
    def setUpClass(cls):
        cls.prs = Presentation('Apresentação1Original.pptx')
        cls.slides = list(cls.prs.slides)
    
    def test_slide_2_has_images(self):
        """Slide 2 should be detected as having visual elements."""
        has_visual = processor.slide_has_visual_elements(self.slides[1])
        self.assertTrue(has_visual, "Slide 2 should have visual elements (images)")
    
    def test_slide_7_has_images(self):
        """Slide 7 should be detected as having visual elements."""
        has_visual = processor.slide_has_visual_elements(self.slides[6])
        self.assertTrue(has_visual, "Slide 7 should have visual elements (images)")
    
    def test_slide_3_no_images(self):
        """Slide 3 should not have visual elements."""
        has_visual = processor.slide_has_visual_elements(self.slides[2])
        self.assertFalse(has_visual, "Slide 3 should not have visual elements")


class TestFontScaling(unittest.TestCase):
    """Test font scaling behavior."""
    
    @classmethod
    def setUpClass(cls):
        cls.prs_orig = Presentation('Apresentação1Original.pptx')
        cls.prs_proc = Presentation('test_output.pptx')
    
    def get_max_font(self, slide):
        """Get maximum font size in slide."""
        max_font = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            max_font = max(max_font, run.font.size.pt)
        return max_font
    
    def test_font_never_shrinks(self):
        """Font size should never be smaller than original."""
        slides_orig = list(self.prs_orig.slides)
        slides_proc = list(self.prs_proc.slides)
        
        for i, (orig, proc) in enumerate(zip(slides_orig, slides_proc)):
            orig_font = self.get_max_font(orig)
            proc_font = self.get_max_font(proc)
            
            if orig_font > 0 and proc_font > 0:
                self.assertGreaterEqual(
                    proc_font,
                    orig_font - 1,  # Allow 1pt tolerance for rounding
                    f"Slide {i+1}: Font shrunk from {orig_font}pt to {proc_font}pt"
                )


class TestTextContentPreservation(unittest.TestCase):
    """Test that text content is preserved during processing."""
    
    @classmethod
    def setUpClass(cls):
        cls.prs_orig = Presentation('Apresentação1Original.pptx')
        cls.prs_proc = Presentation('test_output.pptx')
    
    def get_all_text(self, slide):
        """Get all text from slide, normalized."""
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Normalize whitespace for comparison
                    text = ' '.join(text.split())
                    texts.append(text)
        return texts
    
    def test_text_preserved_all_slides(self):
        """All text content should be preserved after processing."""
        slides_orig = list(self.prs_orig.slides)
        slides_proc = list(self.prs_proc.slides)
        
        for i, (orig, proc) in enumerate(zip(slides_orig, slides_proc)):
            orig_texts = self.get_all_text(orig)
            proc_texts = self.get_all_text(proc)
            
            # Each original text should appear in processed (possibly reformatted)
            for orig_text in orig_texts:
                found = any(orig_text in proc_text or proc_text in orig_text 
                           for proc_text in proc_texts)
                self.assertTrue(
                    found,
                    f"Slide {i+1}: Text content was lost: '{orig_text[:50]}...'"
                )


if __name__ == '__main__':
    # Run tests
    unittest.main(verbosity=2)
