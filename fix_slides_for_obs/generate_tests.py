# -*- coding: utf-8 -*-
"""Script to generate test_individual_slides.py with detailed shape assertions."""
from pptx import Presentation
import fix_slides_for_obs_processor as processor
import os
import sys

# Suppress DEBUG output during generation
class DevNull:
    def write(self, msg): pass
    def flush(self): pass

# Expected fonts for each slide
fonts = {
    2: 36.0, 3: 52.0, 4: 48.0, 6: 67.0, 7: 36.0, 8: 48.0, 9: 48.0, 11: 48.0,
    13: 60.0, 15: 54.0, 17: 85.0, 19: 48.0, 20: 48.0, 21: 48.0, 23: 48.0,
    25: 54.0, 27: 80.0, 29: 54.0, 30: 54.0, 31: 54.0, 33: 81.0, 35: 74.0,
    37: 73.0, 39: 73.0, 41: 48.0, 42: 48.0, 43: 48.0, 44: 48.0, 45: 48.0,
    46: 48.0, 47: 48.0, 49: 47.0, 50: 62.0, 51: 72.0, 52: 51.0, 53: 58.0,
    54: 74.0, 55: 64.0, 56: 74.0, 57: 75.0, 58: 74.0, 59: 72.0, 60: 74.0,
    61: 68.0, 62: 39.0, 64: 39.0, 65: 39.0, 67: 44.0, 68: 44.0, 70: 48.0,
    71: 48.0, 72: 48.0
}

# Visual slides
visual_slides = {2, 7}

# Descriptions
descs = {
    2: 'Has image - should NOT transform',
    3: 'Oração Vocacional',
    4: 'Long prayer text',
    6: 'Mesmo as trevas',
    7: 'Has image - should NOT transform',
    8: 'Que germine o Salvador v1',
    9: 'Que germine o Salvador v2',
    11: 'Senhor e Filho de Deus',
    13: 'Salmo responsorial',
    15: 'Aleluia',
    17: 'Oração da comunidade',
    19: 'Eis, Senhor, a tua vinha v1',
    20: 'Eis, Senhor, a tua vinha v2',
    21: 'Eis, Senhor, a tua vinha v3',
    23: 'Conclusão à Preparação das Oferendas',
    25: 'Santo, Santo, Santo',
    27: 'Enviai o vosso Espírito Santo',
    29: 'Anamnese 1',
    30: 'Anamnese 2',
    31: 'Anamnese 3',
    33: 'Aceitai, ó Senhor, a nossa oferta',
    35: 'O Espírito nos una',
    37: 'Lembrai-vos da Igreja',
    39: 'Concedei-lhes a luz eterna',
    41: 'Cântico de Zacarias v1',
    42: 'Cântico de Zacarias v2',
    43: 'Cântico de Zacarias v3',
    44: 'Cântico de Zacarias v4',
    45: 'Cântico de Zacarias v5',
    46: 'Cântico de Zacarias v6',
    47: 'Cântico de Zacarias v7',
    49: 'Oração a Santo Antônio p1',
    50: 'Oração a Santo Antônio p2',
    51: 'Oração a Santo Antônio p3',
    52: 'Oração a Santo Antônio p4',
    53: 'Responsório de Santo Antônio p1',
    54: 'Responsório refrão',
    55: 'Responsório v2',
    56: 'Responsório refrão',
    57: 'Responsório v3',
    58: 'Responsório refrão',
    59: 'Responsório v4',
    60: 'Responsório refrão',
    61: 'Rogai por nós',
    62: 'Oremos (intercessão)',
    64: 'Benção dos pães',
    65: 'Oremos (pães)',
    67: 'Santo Antônio junto a nós v1-2',
    68: 'Santo Antônio junto a nós v3-4',
    70: 'Hino do Jubileu v1',
    71: 'Hino do Jubileu v2',
    72: 'Hino do Jubileu v3',
}

header = '''# -*- coding: utf-8 -*-
"""
Individual slide tests for fix_slides_for_obs_processor.

Each test:
1. Loads an individual slide file from test_slides/
2. Processes it through the processor
3. Saves to a temporary file
4. Reopens the saved file
5. Verifies the expected values (position, size, font, text content)

Run with: python -m pytest test_individual_slides.py -v
"""
import pytest
import os
from pptx import Presentation
import fix_slides_for_obs_processor as processor


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def get_slide_text(slide):
    """Get all text content from a slide, joined with newlines."""
    all_text = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            all_text.append(shape.text_frame.text.strip())
    return '\\n'.join(all_text)


def get_max_font(slide):
    """Get the maximum font size from all shapes in slide."""
    max_font = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        max_font = max(max_font, run.font.size.pt)
    return max_font if max_font > 0 else None


def get_shapes_info(slide):
    """Get detailed info about each shape with text.
    
    Returns list of dicts with:
    - name: shape name
    - fonts: list of unique font sizes in the shape
    - text: shape text content
    """
    shapes_info = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            fonts_set = set()
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        fonts_set.add(run.font.size.pt)
            shapes_info.append({
                'name': shape.name,
                'fonts': sorted(fonts_set),
                'text': shape.text_frame.text.strip(),
            })
    return shapes_info


def process_and_save_slide(slide_num, tmp_path):
    """Process a slide, save it, and return path to saved file."""
    input_file = f'test_slides/slide_{slide_num:02d}.pptx'
    
    if not os.path.exists(input_file):
        pytest.skip(f"File {input_file} not found")
    
    output_file = tmp_path / f'slide_{slide_num:02d}_processed.pptx'
    
    # Load and process
    prs = Presentation(input_file)
    processor.reposition_and_maximize_font(prs)
    prs.save(str(output_file))
    
    return str(output_file)


def load_and_verify(filepath):
    """Load a processed file and extract info."""
    prs = Presentation(filepath)
    slide = prs.slides[0]
    
    return {
        'text': get_slide_text(slide),
        'max_font': get_max_font(slide),
        'has_visual': processor.slide_has_visual_elements(slide),
        'shapes': get_shapes_info(slide),
        'num_shapes': len([s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]),
    }


# =============================================================================
# TEST FUNCTIONS
# =============================================================================

'''

footer = '''
if __name__ == '__main__':
    pytest.main([__file__, '-v'])
'''

output = header

for slide_num in sorted(fonts.keys()):
    f = os.path.join('test_slides', f'slide_{slide_num:02d}.pptx')
    if os.path.exists(f):
        old_stderr = sys.stderr
        sys.stderr = DevNull()
        try:
            prs = Presentation(f)
            processor.reposition_and_maximize_font(prs)
            slide = prs.slides[0]
            
            # Get text
            all_text = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    all_text.append(shape.text_frame.text.strip())
            full_text = '\n'.join(all_text)
            
            # Get shapes info
            shapes_data = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    shape_fonts = set()
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                shape_fonts.add(run.font.size.pt)
                    shapes_data.append((shape.name, sorted(shape_fonts)))
            
            desc = descs.get(slide_num, f'Slide {slide_num}')
            has_visual = 'True' if slide_num in visual_slides else 'False'
            font = fonts[slide_num]
            num_shapes = len(shapes_data)
            
            output += f'''def test_slide_{slide_num:02d}(tmp_path):
    """{desc}"""
    output = process_and_save_slide({slide_num}, tmp_path)
    result = load_and_verify(output)
    
    assert result['has_visual'] == {has_visual}
    assert result['max_font'] == {font}
    assert result['num_shapes'] == {num_shapes}
    assert result['text'] == {repr(full_text)}
'''
            
            # Add shape assertions
            for i, (name, shape_font_list) in enumerate(shapes_data):
                output += f"    assert result['shapes'][{i}]['name'] == {repr(name)}\n"
                output += f"    assert result['shapes'][{i}]['fonts'] == {shape_font_list}\n"
            
            output += '\n\n'
        finally:
            sys.stderr = old_stderr

output += footer

with open('test_individual_slides.py', 'w', encoding='utf-8') as f:
    f.write(output)

print('Generated test_individual_slides.py with detailed shape assertions')
