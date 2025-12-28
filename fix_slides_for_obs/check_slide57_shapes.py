from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fix_slides_for_obs_processor as processor

# Check if slide 57 has visual elements
prs_orig = Presentation('Apresentação1Original.pptx')

slide = list(prs_orig.slides)[56]  # 0-indexed

print('=== SLIDE 57 SHAPES ===')
for shape in slide.shapes:
    print(f'Shape: {shape.name}, type: {shape.shape_type}')
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        print('  *** THIS IS A PICTURE ***')

print(f'\nslide_has_visual_elements: {processor.slide_has_visual_elements(slide)}')
