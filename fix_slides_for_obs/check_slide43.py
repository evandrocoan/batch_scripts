from pptx import Presentation

# Check slide 43
prs_orig = Presentation('Apresentação1Original.pptx')
prs_proc = Presentation('test_output.pptx')

slide_orig = list(prs_orig.slides)[42]  # 0-indexed
slide_proc = list(prs_proc.slides)[42]

print('=== SLIDE 43 ORIGINAL ===')
for shape in slide_orig.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        print(f'Shape: {shape.name}')
        print(f'  Position: left={shape.left/914400:.2f}" top={shape.top/914400:.2f}"')
        print(f'  Size: width={shape.width/914400:.2f}" height={shape.height/914400:.2f}"')
        text = shape.text_frame.text.strip()
        print(f'  Text: "{text[:60]}..."')
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    size = run.font.size.pt if run.font.size else 'None'
                    print(f'    Font: {size}pt - "{run.text[:40]}"')
        print()

print('\n=== SLIDE 43 AFTER PROCESSING ===')
for shape in slide_proc.shapes:
    if shape.has_text_frame and shape.text_frame.text.strip():
        print(f'Shape: {shape.name}')
        print(f'  Position: left={shape.left/914400:.2f}" top={shape.top/914400:.2f}"')
        print(f'  Size: width={shape.width/914400:.2f}" height={shape.height/914400:.2f}"')
        text = shape.text_frame.text.strip()
        print(f'  Text: "{text[:60]}..."')
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    size = run.font.size.pt if run.font.size else 'None'
                    print(f'    Font: {size}pt - "{run.text[:40]}"')
        print()
