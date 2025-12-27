"""
Shared functions for PowerPoint slide processing with OBS chroma key support.
"""
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor


def apply_solid_glow_to_run(run, color_hex, size_pt):
    """
    Applies a SOLID glow effect (0% transparency) to create a highlighter background effect.
    Uses the effectLst structure that PowerPoint expects.
    
    Args:
        run: The text run to apply the glow to
        color_hex: Hex color code (with or without #)
        size_pt: Glow size in points
    """
    # Strip '#' from color if present (for VS Code color preview support)
    color_hex = color_hex.lstrip('#')
    
    # Get the Run Properties (rPr) element or create it
    rPr = run._r.get_or_add_rPr()
    
    # Calculate radius in EMUs (1 point = 12700 EMUs)
    radius_emu = int(size_pt * 12700)
    
    # Create the XML for the effect list with glow
    # Using alpha val="100000" for 100% opacity (0% transparency)
    effectlst_xml = f'''<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:glow rad="{radius_emu}">
            <a:srgbClr val="{color_hex}">
                <a:alpha val="100000"/>
            </a:srgbClr>
        </a:glow>
    </a:effectLst>'''
    
    # Remove any existing effectLst or standalone glow
    existing_effectlst = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
    if existing_effectlst is not None:
        rPr.remove(existing_effectlst)
    
    existing_glow = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}glow")
    if existing_glow is not None:
        rPr.remove(existing_glow)
    
    # Parse the new XML and inject it into the run properties
    # effectLst should come after solidFill but before font typefaces
    effectlst_element = parse_xml(effectlst_xml)
    
    # Insert after solidFill if it exists, otherwise at the beginning
    solid_fill = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    if solid_fill is not None:
        # Insert after solidFill
        insert_index = list(rPr).index(solid_fill) + 1
        rPr.insert(insert_index, effectlst_element)
    else:
        rPr.insert(0, effectlst_element)


def process_presentation(prs, glow_color, glow_size, text_color):
    """
    Process a PowerPoint presentation to add glow effects and set backgrounds.
    
    Args:
        prs: PowerPoint Presentation object
        glow_color: Hex color for glow (with or without #)
        glow_size: Glow size in points
        text_color: Hex color for text (with or without #)
    
    Returns:
        int: Number of text shapes processed
    """
    # Strip '#' from text color if present
    text_color = text_color.lstrip('#')
    
    count = 0
    
    # Iterate through every slide
    for slide in prs.slides:
        # Check if slide has any text
        has_text = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                has_text = True
                break
        
        # Set slide background: white if has text, black if no text
        background = slide.background
        fill = background.fill
        fill.solid()
        if has_text:
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        else:
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
            continue  # Skip processing shapes on empty slides
        
        # Iterate through every shape on the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if not shape.has_text_frame:
                continue
            
            # Process all text regardless of whether it's from master or manual text box
            text_processed = False
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():  # Only process non-empty runs
                        # Apply the SOLID glow (0% transparency) to create highlighter effect
                        apply_solid_glow_to_run(run, glow_color, glow_size)
                        
                        # Set text color
                        run.font.color.rgb = RGBColor(
                            int(text_color[0:2], 16),
                            int(text_color[2:4], 16),
                            int(text_color[4:6], 16)
                        )
                        text_processed = True
            
            if text_processed:
                count += 1
    
    return count
