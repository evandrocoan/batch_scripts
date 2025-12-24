from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ================= CONFIGURATION =================
INPUT_FILE = "Apresentação1.pptx"  # Put your file name here
OUTPUT_FILE = "Apresentação1_fixed.pptx"
GLOW_COLOR = "#FFFFF0"  # Lighter, more discrete yellow
GLOW_SIZE_PT = 18      # Size of the glow in points (reduce to avoid overlap between letters)
TEXT_COLOR = "#010101"  # Hex code for near-black text (R=1, G=1, B=1)
BACKGROUND_COLOR = "#000000"  # Hex code for black background (R=0, G=0, B=0)
# =================================================

def apply_solid_glow_to_run(run, color_hex, size_pt):
    """
    Applies a SOLID glow effect (0% transparency) to create a highlighter background effect.
    Uses the effectLst structure that PowerPoint expects.
    """
    # Strip '#' from color if present (for VS Code color preview support)
    color_hex = color_hex.lstrip('#')
    
    # 1. Get the Run Properties (rPr) element or create it
    rPr = run._r.get_or_add_rPr()

    # 2. Calculate radius in EMUs (1 point = 12700 EMUs)
    radius_emu = int(size_pt * 12700)

    # 3. Create the XML for the effect list with glow
    # Using alpha val="100000" for 100% opacity (0% transparency)
    effectlst_xml = f'''<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:glow rad="{radius_emu}">
            <a:srgbClr val="{color_hex}">
                <a:alpha val="100000"/>
            </a:srgbClr>
        </a:glow>
    </a:effectLst>'''
    
    # 4. Remove any existing effectLst or standalone glow
    existing_effectlst = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst")
    if existing_effectlst is not None:
        rPr.remove(existing_effectlst)
    
    existing_glow = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}glow")
    if existing_glow is not None:
        rPr.remove(existing_glow)

    # 5. Parse the new XML and inject it into the run properties
    # effectLst should come after solidFill but before font typefaces
    effectlst_element = parse_xml(effectlst_xml)
    
    # Insert after solidFill if it exists, otherwise at the beginning
    solid_fill = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    if solid_fill is not None:
        # Insert after solidFill
        insert_index = list(rPr).index(solid_fill) + 1
        rPr.insert(insert_index, effectlst_element)
    else:
        rPr.insert(0, effectlst_element)

def process_presentation():
    print(f"Opening {INPUT_FILE}...")
    prs = Presentation(INPUT_FILE)
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    count = 0

    # Iterate through every slide
    for slide in prs.slides:
        # Check if slide has any text
        has_text = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                has_text = True
                break
        
        # Set slide background: white if has text, black if no text
        background = slide.background
        fill = background.fill
        fill.solid()
        if has_text:
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        else:
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
        
        # Iterate through every shape on the slide
        for shape in slide.shapes:
            # Check if the shape has text
            if not shape.has_text_frame:
                continue

            # Process all text regardless of whether it's from master or manual text box
            text_processed = False
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():  # Only process non-empty runs
                        # Apply the SOLID glow (0% transparency) to create highlighter effect
                        apply_solid_glow_to_run(run, GLOW_COLOR, GLOW_SIZE_PT)
                        
                        # Set text color to near-black (#010101)
                        run.font.color.rgb = RGBColor(1, 1, 1)
                        text_processed = True
            
            if text_processed:
                count += 1
                shape_type = "Placeholder" if shape.is_placeholder else "TextBox"
                print(f"  Applied glow to {shape_type}: {shape.name} - '{shape.text_frame.text[:50]}'...")

    print(f"Processed {count} text shapes.")
    print(f"Saving to {OUTPUT_FILE}...")
    prs.save(OUTPUT_FILE)
    print("Done!")

if __name__ == "__main__":
    try:
        process_presentation()
    except FileNotFoundError:
        print(f"Error: Could not find file '{INPUT_FILE}'. Please check the name.")
