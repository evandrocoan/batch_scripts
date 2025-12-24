from pptx import Presentation
import zipfile
import xml.etree.ElementTree as ET

# This script will extract and show the XML structure of text with glow
INPUT_FILE = "Apresentação1_fixed.pptx"

def inspect_presentation():
    """
    Extracts the raw XML to see how PowerPoint stores the glow effect
    """
    print(f"Inspecting {INPUT_FILE}...")
    
    # Method 1: Using python-pptx to inspect runs
    prs = Presentation(INPUT_FILE)
    
    print("\n=== METHOD 1: Inspecting via python-pptx ===")
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                print(f"\nShape: {shape.name} (Type: {'Placeholder' if shape.is_placeholder else 'TextBox'})")
                print(f"Text preview: {shape.text_frame.text[:50]}...")
                
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            print(f"  Run {run_idx}: '{run.text[:30]}...'")
                            # Get the raw XML of this run
                            run_xml = ET.tostring(run._r, encoding='unicode')
                            print(f"  XML (first 500 chars): {run_xml[:500]}")
                            if 'glow' in run_xml.lower():
                                print(f"  *** FOUND GLOW in this run! ***")
                                print(f"  Full XML:\n{run_xml}")
                            break  # Only check first run per paragraph
                    break  # Only check first paragraph per shape
        
        if slide_num >= 2:  # Only check first 2 slides
            break
    
    # Method 2: Direct ZIP inspection
    print("\n\n=== METHOD 2: Direct ZIP inspection ===")
    print("Looking at slide1.xml directly...")
    
    with zipfile.ZipFile(INPUT_FILE, 'r') as zip_ref:
        # Read the first slide's XML
        slide1_xml = zip_ref.read('ppt/slides/slide1.xml').decode('utf-8')
        
        # Look for glow elements
        if 'glow' in slide1_xml.lower():
            print("Found 'glow' in slide1.xml!")
            # Extract just the sections with glow
            lines = slide1_xml.split('>')
            for i, line in enumerate(lines):
                if 'glow' in line.lower():
                    # Print context around the glow
                    start = max(0, i-3)
                    end = min(len(lines), i+4)
                    context = '>'.join(lines[start:end])
                    print(f"\nGlow context:\n{context}\n")
        else:
            print("No 'glow' found in slide1.xml")
            print(f"First 1000 chars of slide1.xml:\n{slide1_xml[:1000]}")

if __name__ == "__main__":
    try:
        inspect_presentation()
    except FileNotFoundError:
        print(f"Error: Could not find file '{INPUT_FILE}'")
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
